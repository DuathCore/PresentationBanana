#!/usr/bin/env python3
"""
PresentationBanana — Academic Presentation Builder v5
======================================================
Uses Excalidraw diagram PNGs for content slides,
Imagen PNGs for title/section/closing.
"""
import json, sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print(json.dumps({"ok": False, "error": "pip install python-pptx"}))
    sys.exit(1)

# ── Dimensions ────────────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL = 19

# ── Colors ────────────────────────────────────────────────────────────────
C_BG      = RGBColor(0x12, 0x1A, 0x2E)
C_BG_DARK = RGBColor(0x0D, 0x12, 0x20)
C_ACCENT  = RGBColor(0xF0, 0xAB, 0x00)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT   = RGBColor(0xC8, 0xD6, 0xE5)
C_MUTED   = RGBColor(0x7A, 0x8B, 0xA0)

# ── Paths ─────────────────────────────────────────────────────────────────
IMG_DIR = Path(__file__).parent.parent / "output" / "images"
SLUG = "from-prompts-to-populations"

def img_diagram(n):
    return IMG_DIR / f"{SLUG}_s{n:02d}_diagram_v5.png"

def img_imagen(n, stype):
    return IMG_DIR / f"{SLUG}_s{n:02d}_{stype}.png"


# ── Helpers ───────────────────────────────────────────────────────────────

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _bg(slide, color=None):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color or C_BG

def _rect(slide, left, top, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def _text(slide, text, left, top, w, h, size=20, bold=False, color=C_WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box

def _notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text

def _img(slide, path, left, top, w, h):
    if path and Path(path).exists():
        slide.shapes.add_picture(str(path), left, top, w, h)

def _footer(slide, num):
    _rect(slide, Inches(0), Inches(7.25), SLIDE_W, Inches(0.04), C_ACCENT)
    _text(slide, f"{num} / {TOTAL}", Inches(12.0), Inches(7.28), Inches(1.2), Inches(0.22),
          size=9, color=C_MUTED, align=PP_ALIGN.RIGHT)
    _text(slide, "Baraghith (2025) · From Prompts to Populations",
          Inches(0.3), Inches(7.28), Inches(5), Inches(0.22), size=9, color=C_MUTED)

def _title_bar(slide, title):
    _text(slide, title, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
          size=28, bold=True, color=C_WHITE)
    _rect(slide, Inches(0.5), Inches(1.05), Inches(6.0), Inches(0.04), C_ACCENT)


# ── Diagram slide helper ─────────────────────────────────────────────────

def _diagram_slide(prs, num, title, diagram_path, notes_text):
    """Standard content slide: title bar + full-width Excalidraw diagram PNG."""
    slide = _blank(prs)
    _bg(slide)
    _title_bar(slide, title)
    # Diagram image: full width below title, above footer
    _img(slide, diagram_path, Inches(0.3), Inches(1.3), Inches(12.7), Inches(5.8))
    _footer(slide, num)
    _notes(slide, notes_text)
    return slide


# ══════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    slide = _blank(prs)
    _bg(slide)
    _img(slide, img_imagen(1, "title"), Inches(7.4), Inches(0.3), Inches(5.6), Inches(6.9))
    _rect(slide, Inches(7.1), Inches(0.3), Inches(0.06), Inches(6.9), C_ACCENT)
    _text(slide, "From Prompts\nto Populations", Inches(0.5), Inches(1.2), Inches(6.4), Inches(2.5),
          size=40, bold=True)
    _rect(slide, Inches(0.5), Inches(3.8), Inches(4.0), Inches(0.06), C_ACCENT)
    _text(slide, "Cultural Evolution as a Meta-Frame\nfor AI Ontology and Individuation",
          Inches(0.5), Inches(4.1), Inches(6.4), Inches(1.5), size=20, color=C_LIGHT)
    _text(slide, "Baraghith (2025)", Inches(0.5), Inches(5.5), Inches(6), Inches(0.4),
          size=16, color=C_ACCENT)
    _notes(slide, "Based on Baraghith (2025). The paper asks: How many kinds of AI are out there?")

def slide_02(prs):
    _diagram_slide(prs, 2, "AI individuation depends on who you ask",
                   img_diagram(2),
                   "Different communities individuate AI differently. Engineers version by weights, RL researchers track policies, auditors need instance IDs.")

def slide_03(prs):
    _diagram_slide(prs, 3, "These six criteria systematically conflict",
                   img_diagram(3),
                   "Same model + different prompts = same under model-instance but different under behavioural. These conflicts are structural.")

def slide_04_section(prs):
    slide = _blank(prs)
    _bg(slide)
    _img(slide, img_imagen(3, "section"), Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    _rect(slide, Inches(0), Inches(5.2), SLIDE_W, Inches(2.3), C_BG_DARK)
    _rect(slide, Inches(0), Inches(5.2), SLIDE_W, Inches(0.04), C_ACCENT)
    _text(slide, "The Ontology Problem", Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.0),
          size=34, bold=True)
    _text(slide, 'Four philosophical positions on what AI "is"',
          Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.9), size=18, color=C_LIGHT)
    _footer(slide, 4)
    _notes(slide, "Before the evolutionary argument, the paper surveys four recent philosophical positions.")

def slide_05(prs):
    _diagram_slide(prs, 5, "Four positions span individuation and agency",
                   img_diagram(5),
                   "These four positions differ on how strongly they reify AI as individuals and where they locate agency.")

def slide_06(prs):
    _diagram_slide(prs, 6, "No single criterion resolves AI identity",
                   img_diagram(6),
                   "Pluralism merely tolerates multiple criteria. A structural explanation is needed — this is where CET enters.")

def slide_07_section(prs):
    slide = _blank(prs)
    _bg(slide)
    _img(slide, img_imagen(7, "section"), Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    _rect(slide, Inches(0), Inches(5.2), SLIDE_W, Inches(2.3), C_BG_DARK)
    _rect(slide, Inches(0), Inches(5.2), SLIDE_W, Inches(0.04), C_ACCENT)
    _text(slide, "Cultural Evolutionary Theory", Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.0),
          size=34, bold=True)
    _text(slide, "A unifying meta-frame, not just pluralism",
          Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.9), size=18, color=C_LIGHT)
    _footer(slide, 7)
    _notes(slide, "CET offers more than tolerance of multiple criteria. It provides a structural explanation.")

def slide_08(prs):
    _diagram_slide(prs, 8, "CET explains what pluralism merely tolerates",
                   img_diagram(8),
                   "Pluralism says 'use whichever criterion fits'. CET says 'each criterion targets a specific evolutionary level'.")

def slide_09(prs):
    _diagram_slide(prs, 9, "AI development exhibits variation, inheritance, selection",
                   img_diagram(9),
                   "CET requires three conditions: variation, inheritance, and selection. AI development fulfils all three.")

def slide_10_section(prs):
    slide = _blank(prs)
    _bg(slide)
    _img(slide, img_imagen(9, "section"), Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    _rect(slide, Inches(0), Inches(5.2), SLIDE_W, Inches(2.3), C_BG_DARK)
    _rect(slide, Inches(0), Inches(5.2), SLIDE_W, Inches(0.04), C_ACCENT)
    _text(slide, "The Four-Level Hierarchy", Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.0),
          size=34, bold=True)
    _text(slide, "Models — Systems — Episodes — Outputs",
          Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.9), size=18, color=C_LIGHT)
    _footer(slide, 10)
    _notes(slide, "The core proposal: AI development operates across four nested levels.")

def slide_11(prs):
    _diagram_slide(prs, 11, "Four nested levels organise AI evolution",
                   img_diagram(11),
                   "Level 3 = weight sets. Level 2 = model+prompts+tools. Level 1 = chat sessions. Level 0 = outputs.")

def slide_12(prs):
    _diagram_slide(prs, 12, "Each level has distinct evolutionary dynamics",
                   img_diagram(12),
                   "Each level has its own variation, transmission, and selection mechanisms.")

def slide_13(prs):
    _diagram_slide(prs, 13, "Downward causation constrains; upward aggregates",
                   img_diagram(13),
                   "Higher levels constrain lower-level variation. Lower levels aggregate into selection signals. L0→L3 bypass is crucial.")

def slide_14(prs):
    _diagram_slide(prs, 14, "Criteria map onto hierarchy levels, not against each other",
                   img_diagram(14),
                   "The six criteria are not rivals. Each targets a specific level.")

def slide_15(prs):
    _diagram_slide(prs, 15, "RLHF demonstrates multi-level evolutionary dynamics",
                   img_diagram(15),
                   "RLHF demonstrates multi-level dynamics: episodes generate selection signals, configs filter, weights reshape, outputs circulate.")

def slide_16(prs):
    _diagram_slide(prs, 16, "Governance interventions cascade through levels",
                   img_diagram(16),
                   "Mandatory watermarks (L0) reshape which configs remain competitive. Restricting training data (L3) constrains everything downstream.")

def slide_17(prs):
    _diagram_slide(prs, 17, "CET separates descriptive from normative questions",
                   img_diagram(17),
                   "CET is independent of consciousness debates. It grounds individuation in operational mechanisms.")

def slide_18(prs):
    _diagram_slide(prs, 18, "Three questions CET must still answer",
                   img_diagram(18),
                   "The framework needs: reproduction predicates, fitness notions, and governance pressure-testing.")

def slide_19_closing(prs):
    slide = _blank(prs)
    _bg(slide)
    _img(slide, img_imagen(18, "closing"), Inches(0), Inches(0), Inches(5.8), SLIDE_H)
    _rect(slide, Inches(5.8), Inches(0.3), Inches(0.06), Inches(6.9), C_ACCENT)
    _text(slide, "One Process,\nMany Levels", Inches(6.3), Inches(1.5), Inches(6.5), Inches(2.5),
          size=36, bold=True)
    _rect(slide, Inches(6.3), Inches(3.8), Inches(3.0), Inches(0.06), C_ACCENT)
    _text(slide, "AI individuation is not a choice\nbetween criteria — it\u2019s recognising\nwhere each belongs in the\nevolutionary hierarchy.",
          Inches(6.3), Inches(4.1), Inches(6.5), Inches(2.0), size=20, color=C_LIGHT)
    _footer(slide, 19)
    _notes(slide, "Take-home: CET provides a meta-frame where all six criteria are partial perspectives on a single multi-level evolutionary process.")


# ══════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════

def main():
    project_dir = Path(__file__).parent.parent
    out_dir = project_dir / "output" / "presentations"
    out_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02(prs)
    slide_03(prs)
    slide_04_section(prs)
    slide_05(prs)
    slide_06(prs)
    slide_07_section(prs)
    slide_08(prs)
    slide_09(prs)
    slide_10_section(prs)
    slide_11(prs)
    slide_12(prs)
    slide_13(prs)
    slide_14(prs)
    slide_15(prs)
    slide_16(prs)
    slide_17(prs)
    slide_18(prs)
    slide_19_closing(prs)

    out_path = out_dir / "from-prompts-to-populations_v5.pptx"
    prs.save(str(out_path))

    print(json.dumps({
        "ok": True,
        "path": str(out_path),
        "slides": TOTAL,
        "version": 5,
    }, indent=2))


if __name__ == "__main__":
    main()
