#!/usr/bin/env python3
"""
PresentationBanana — PowerPoint Builder
========================================
Reads workspace/slide_structure.md and output/images/v{N}/ to
build a professional .pptx presentation.

Usage:
    python build_pptx.py --version 1

Output:
    JSON to stdout: {"ok": true, "path": "...", "slides": 10}
    PPTX saved to: ../output/presentations/presentation_v{N}.pptx

Requirements:
    pip install python-pptx pillow
"""

import argparse
import json
import re
import sys
from pathlib import Path


def _get_project_slug(workspace: Path) -> str:
    """Derive a URL-safe slug from workspace/slide_structure.md topic."""
    structure_file = workspace / "slide_structure.md"
    try:
        content = structure_file.read_text(encoding="utf-8")
        m = re.search(r'\*\*Topic:\*\*\s*(.+)', content)
        if not m:
            return "presentation"
        topic = m.group(1).strip()
        topic = re.split(r'[—–,:]', topic)[0].strip()
        for old, new in [('ä','ae'),('ö','oe'),('ü','ue'),('Ä','Ae'),('Ö','Oe'),('Ü','Ue'),('ß','ss')]:
            topic = topic.replace(old, new)
        words = topic.lower().split()[:4]
        slug = '-'.join(re.sub(r'[^a-z0-9]', '', w) for w in words)
        slug = re.sub(r'-+', '-', slug).strip('-')
        return slug or "presentation"
    except Exception:
        return "presentation"


def _find_image(images_dir: Path, slug: str, num: int, slide_type: str) -> Path:
    """Resolve image path — new slug-based naming with fallback to old slide_N.png."""
    new_path = images_dir / f"{slug}_s{num:02d}_{slide_type}.png"
    if new_path.exists():
        return new_path
    # Fallback for images generated before slug naming
    old_path = images_dir / f"slide_{num}.png"
    return old_path


def _find_icon(images_dir: Path, slug: str, num: int, icon_idx: int) -> Path:
    """Resolve icon path with fallback."""
    new_path = images_dir / f"{slug}_s{num:02d}_icon_{icon_idx}.png"
    if new_path.exists():
        return new_path
    return images_dir / f"slide_{num}_icon_{icon_idx}.png"

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
except ImportError:
    print(json.dumps({"ok": False, "error": "python-pptx not installed. Run: pip install python-pptx"}))
    sys.exit(1)

# ── Slide dimensions (16:9 widescreen) ────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Color palettes (selectable via --style) ────────────────────────────────────
STYLES = {
    "dark-professional": {   # default: deep navy + gold
        "bg":          RGBColor(0x12, 0x1A, 0x2E),
        "bg_mid":      RGBColor(0x1A, 0x27, 0x44),
        "accent":      RGBColor(0xF0, 0xAB, 0x00),
        "white":       RGBColor(0xFF, 0xFF, 0xFF),
        "light":       RGBColor(0xC8, 0xD6, 0xE5),
        "placeholder": RGBColor(0x1E, 0x30, 0x50),
    },
    "light-modern": {        # white background, navy text, teal accent
        "bg":          RGBColor(0xF8, 0xF9, 0xFA),
        "bg_mid":      RGBColor(0xE9, 0xEC, 0xF0),
        "accent":      RGBColor(0x00, 0x87, 0x8A),
        "white":       RGBColor(0x1A, 0x1A, 0x2E),   # inverted (dark text)
        "light":       RGBColor(0x44, 0x55, 0x66),
        "placeholder": RGBColor(0xCC, 0xD6, 0xE0),
    },
    "minimal": {             # pure white, black, subtle grey
        "bg":          RGBColor(0xFF, 0xFF, 0xFF),
        "bg_mid":      RGBColor(0xF0, 0xF0, 0xF0),
        "accent":      RGBColor(0x33, 0x33, 0x33),
        "white":       RGBColor(0x11, 0x11, 0x11),
        "light":       RGBColor(0x55, 0x55, 0x55),
        "placeholder": RGBColor(0xDD, 0xDD, 0xDD),
    },
    "bold-creative": {       # near-black, vibrant orange-red
        "bg":          RGBColor(0x0D, 0x0D, 0x0D),
        "bg_mid":      RGBColor(0x1A, 0x1A, 0x1A),
        "accent":      RGBColor(0xFF, 0x45, 0x00),
        "white":       RGBColor(0xFF, 0xFF, 0xFF),
        "light":       RGBColor(0xBB, 0xBB, 0xBB),
        "placeholder": RGBColor(0x2A, 0x2A, 0x2A),
    },
    "dhl-corporate": {       # DHL brand: dark charcoal + DHL yellow
        "bg":          RGBColor(0x1A, 0x1A, 0x1A),
        "bg_mid":      RGBColor(0x26, 0x26, 0x26),
        "accent":      RGBColor(0xFF, 0xCC, 0x00),
        "white":       RGBColor(0xFF, 0xFF, 0xFF),
        "light":       RGBColor(0xCC, 0xCC, 0xCC),
        "placeholder": RGBColor(0x32, 0x32, 0x32),
    },
}

# Active palette (set in main())
_P = STYLES["dark-professional"]

# Shorthand references (updated in main() via set_style())
C_BG          = _P["bg"]
C_BG_MID      = _P["bg_mid"]
C_ACCENT      = _P["accent"]
C_WHITE       = _P["white"]
C_LIGHT       = _P["light"]
C_PLACEHOLDER = _P["placeholder"]


def set_style(name: str):
    """Apply a named style to all color constants."""
    global C_BG, C_BG_MID, C_ACCENT, C_WHITE, C_LIGHT, C_PLACEHOLDER
    palette = STYLES.get(name, STYLES["dark-professional"])
    C_BG          = palette["bg"]
    C_BG_MID      = palette["bg_mid"]
    C_ACCENT      = palette["accent"]
    C_WHITE       = palette["white"]
    C_LIGHT       = palette["light"]
    C_PLACEHOLDER = palette["placeholder"]


# ── Markdown parser ────────────────────────────────────────────────────────────

def parse_slide_structure(workspace_path: Path) -> dict:
    """Parse workspace/slide_structure.md into a dict of meta + slides list."""
    md_file = workspace_path / "slide_structure.md"
    if not md_file.exists():
        raise FileNotFoundError(f"Not found: {md_file}")

    content = md_file.read_text(encoding="utf-8")
    sections = re.split(r"\n---\n", content)

    # Meta header (first section)
    meta = {}
    for line in sections[0].splitlines():
        m = re.match(r"\*\*([\w ]+):\*\*\s*(.+)", line)
        if m:
            key = m.group(1).strip().lower().replace(" ", "_")
            meta[key] = m.group(2).strip()

    # Individual slides
    slides = []
    for section in sections[1:]:
        slide = _parse_slide_section(section.strip())
        if slide:
            slides.append(slide)

    return {"meta": meta, "slides": slides}


def _parse_slide_section(text: str) -> dict | None:
    if not text:
        return None

    def get(pattern, default=""):
        m = re.search(pattern, text)
        return m.group(1).strip() if m else default

    # Slide number
    num_m = re.search(r"## Slide (\d+)", text)
    num = int(num_m.group(1)) if num_m else None

    slide_type = get(r"\*\*Type:\*\*\s*(\w+)", "content").lower()
    title       = get(r"\*\*Title:\*\*\s*(.+)")
    subtitle    = get(r"\*\*Subtitle:\*\*\s*(.+)")
    speaker     = get(r"\*\*Speaker Notes:\*\*\s*(.+)")
    img_concept = get(r"\*\*Image Concept:\*\*\s*(.+)")

    # Bullets (multi-line block after **Bullets:**)
    bullets = []
    bm = re.search(r"\*\*Bullets:\*\*\n((?:[-*] .+\n?)+)", text)
    if bm:
        for line in bm.group(1).splitlines():
            clean = re.sub(r"^[-*]\s*", "", line).strip()
            if clean:
                bullets.append(clean)

    if not title:
        return None

    has_image  = get(r"\*\*Image:\*\*\s*(\w+)", "yes").lower()
    image_type = get(r"\*\*Image Type:\*\*\s*(\w+)", slide_type).lower()

    # Icon concepts (up to 3, listed as "Icon 1:", "Icon 2:", "Icon 3:")
    icons = []
    for idx in range(1, 4):
        ic = get(rf"\*\*Icon {idx}:\*\*\s*(.+)", "")
        if ic and ic.lower() != "none":
            icons.append(ic)

    return {
        "number":        num,
        "type":          slide_type,
        "has_image":     has_image not in ("no", "none", "false"),
        "image_type":    image_type,
        "title":         title,
        "subtitle":      subtitle,
        "bullets":       bullets,
        "icons":         icons,
        "speaker_notes": speaker,
        "image_concept": img_concept,
    }


# ── Slide building helpers ─────────────────────────────────────────────────────

def _blank_slide(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[6])  # blank layout


def _bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, left, top, width, height, color: RGBColor):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _textbox(slide, text: str, left, top, width, height,
             size=18, bold=False, color=C_WHITE, align=PP_ALIGN.LEFT,
             wrap=True):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def _image_or_placeholder(slide, path: Path, left, top, width, height):
    if path and path.exists():
        slide.shapes.add_picture(str(path), left, top, width, height)
    else:
        shape = _rect(slide, left, top, width, height, C_PLACEHOLDER)
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "[Image not generated yet]"
        run.font.size = Pt(11)
        run.font.color.rgb = C_LIGHT


def _speaker_notes(slide, text: str):
    if text:
        slide.notes_slide.notes_text_frame.text = text


# ── Slide type builders ────────────────────────────────────────────────────────

def build_title_slide(prs, data: dict, img: Path):
    """
    Layout:
        Left  (0 → 7.2"): Title (large) + accent line + Subtitle
        Right (7.5" → end): Image
    """
    slide = _blank_slide(prs)
    _bg(slide, C_BG)

    # Image — right panel
    _image_or_placeholder(slide, img, Inches(7.4), Inches(0.3), Inches(5.6), Inches(6.9))

    # Vertical accent bar between text and image
    _rect(slide, Inches(7.1), Inches(0.3), Inches(0.05), Inches(6.9), C_ACCENT)

    # Title
    _textbox(slide, data["title"], Inches(0.5), Inches(1.8), Inches(6.4), Inches(2.2),
             size=38, bold=True, color=C_WHITE)

    # Gold accent line
    _rect(slide, Inches(0.5), Inches(4.15), Inches(3.5), Inches(0.06), C_ACCENT)

    # Subtitle
    if data.get("subtitle"):
        _textbox(slide, data["subtitle"], Inches(0.5), Inches(4.4), Inches(6.4), Inches(1.5),
                 size=20, color=C_LIGHT)

    _speaker_notes(slide, data.get("speaker_notes", ""))
    return slide


def build_content_slide(prs, data: dict, img: Path):
    """
    Layout:
        Left  (0 → 7.2"): Title bar + bullet points
        Right (7.5" → end): Image (full height)
    """
    slide = _blank_slide(prs)
    _bg(slide, C_BG)

    TEXT_W  = Inches(7.2)
    IMG_L   = Inches(7.5)
    IMG_W   = Inches(5.5)
    MARGIN  = Inches(0.45)

    # Image — right panel
    _image_or_placeholder(slide, img, IMG_L, Inches(0.0), IMG_W, SLIDE_H)

    # Title
    _textbox(slide, data["title"], MARGIN, Inches(0.3), TEXT_W - MARGIN, Inches(0.85),
             size=26, bold=True, color=C_WHITE)

    # Gold accent line under title
    _rect(slide, MARGIN, Inches(1.18), Inches(4.2), Inches(0.05), C_ACCENT)

    # Bullets
    bullets = data.get("bullets", [])
    if bullets:
        box = slide.shapes.add_textbox(MARGIN, Inches(1.35), TEXT_W - MARGIN - Inches(0.1), Inches(5.8))
        tf = box.text_frame
        tf.word_wrap = True
        for i, bullet in enumerate(bullets):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(10)
            run = p.add_run()
            run.text = f"\u25cf  {bullet}"   # ● bullet character
            run.font.size = Pt(17)
            run.font.color.rgb = C_LIGHT

    # Right-edge thin accent strip
    _rect(slide, IMG_L - Inches(0.06), Inches(0.0), Inches(0.06), SLIDE_H, C_ACCENT)

    _speaker_notes(slide, data.get("speaker_notes", ""))
    return slide


def build_section_slide(prs, data: dict, img: Path = None):
    """
    Full-width section header.
    If an image is provided: image fills slide as background with dark overlay + title.
    Without image: solid colored background with centered title.
    """
    slide = _blank_slide(prs)

    if img and img.exists():
        # Full-bleed background image
        _bg(slide, C_BG)
        slide.shapes.add_picture(str(img), Inches(0.0), Inches(0.0), SLIDE_W, SLIDE_H)

        # Narrow dark overlay strip at bottom for text legibility
        _rect(slide, Inches(0.0), Inches(5.2), SLIDE_W, Inches(2.3),
              RGBColor(0x00, 0x00, 0x00))  # compact black stripe

        # Horizontal accent line above title
        _rect(slide, Inches(0.0), Inches(5.2), SLIDE_W, Inches(0.04), C_ACCENT)

        # Title over image
        _textbox(slide, data["title"],
                 Inches(0.8), Inches(5.35), Inches(11.7), Inches(1.0),
                 size=34, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        if data.get("subtitle"):
            _textbox(slide, data["subtitle"],
                     Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.9),
                     size=18, color=C_LIGHT, align=PP_ALIGN.CENTER)
    else:
        # Fallback: solid background, centered
        _bg(slide, C_BG_MID)
        _rect(slide, Inches(0.0), Inches(2.8), SLIDE_W, Inches(0.04), C_ACCENT)

        _textbox(slide, data["title"],
                 Inches(0.8), Inches(2.2), Inches(11.7), Inches(1.8),
                 size=34, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        if data.get("subtitle"):
            _textbox(slide, data["subtitle"],
                     Inches(0.8), Inches(3.8), Inches(11.7), Inches(1.2),
                     size=20, color=C_LIGHT, align=PP_ALIGN.CENTER)

        _rect(slide, Inches(0.0), Inches(4.7), SLIDE_W, Inches(0.04), C_ACCENT)

    _speaker_notes(slide, data.get("speaker_notes", ""))
    return slide


def build_visual_slide(prs, data: dict, img: Path):
    """
    Visual explanation slide — image IS the content.
    Layout:
        Top:    Small title + accent line
        Center: Large image (nearly full-width), no bullet text
        Bottom: Optional subtitle as caption
    Use when an illustration, diagram, or visual metaphor replaces a bullet list.
    """
    slide = _blank_slide(prs)
    _bg(slide, C_BG)

    MARGIN = Inches(0.4)

    # Title (smaller than content slides)
    _textbox(slide, data["title"], MARGIN, Inches(0.2), SLIDE_W - MARGIN * 2, Inches(0.75),
             size=24, bold=True, color=C_WHITE)

    # Accent line under title
    _rect(slide, MARGIN, Inches(0.98), SLIDE_W - MARGIN * 2, Inches(0.04), C_ACCENT)

    # Large centered image — primary content
    _image_or_placeholder(slide, img, MARGIN, Inches(1.1), SLIDE_W - MARGIN * 2, Inches(5.5))

    # Optional caption / subtitle at bottom
    if data.get("subtitle"):
        _textbox(slide, data["subtitle"],
                 MARGIN, Inches(6.7), SLIDE_W - MARGIN * 2, Inches(0.65),
                 size=12, color=C_LIGHT, align=PP_ALIGN.CENTER)

    _speaker_notes(slide, data.get("speaker_notes", ""))
    return slide


def build_closing_slide(prs, data: dict, img: Path):
    """Same structure as title slide, adapted for closing."""
    slide = _blank_slide(prs)
    _bg(slide, C_BG)

    # Image — left panel for variety
    _image_or_placeholder(slide, img, Inches(0.0), Inches(0.0), Inches(5.8), SLIDE_H)

    # Gold separator
    _rect(slide, Inches(5.8), Inches(0.3), Inches(0.05), Inches(6.9), C_ACCENT)

    # Title
    _textbox(slide, data["title"],
             Inches(6.2), Inches(2.0), Inches(6.8), Inches(2.0),
             size=34, bold=True, color=C_WHITE)

    _rect(slide, Inches(6.2), Inches(4.1), Inches(3.0), Inches(0.06), C_ACCENT)

    # Subtitle / CTA
    if data.get("subtitle"):
        _textbox(slide, data["subtitle"],
                 Inches(6.2), Inches(4.35), Inches(6.8), Inches(1.5),
                 size=18, color=C_LIGHT)

    _speaker_notes(slide, data.get("speaker_notes", ""))
    return slide


def build_data_slide(prs, data: dict, images_dir: Path, slug: str = "presentation"):
    """
    Data/agenda slide — no side image, clean layout.
    If icons exist (slide_N_icon_1/2/3.png), shows them in a row with bullet text below.
    Otherwise: clean text-only with large accent number/stat support.
    """
    slide = _blank_slide(prs)
    _bg(slide, C_BG)

    MARGIN = Inches(0.6)

    # Title
    _textbox(slide, data["title"], MARGIN, Inches(0.3), SLIDE_W - MARGIN * 2, Inches(0.85),
             size=26, bold=True, color=C_WHITE)
    _rect(slide, MARGIN, Inches(1.18), Inches(4.0), Inches(0.04), C_ACCENT)

    num     = data.get("number", 1)
    icons   = data.get("icons", [])
    bullets = data.get("bullets", [])

    # ── Layout A: Icons vorhanden → Icon-Reihe ──────────────────────────────
    if icons:
        icon_count = len(icons)
        icon_size  = Inches(1.6)
        total_w    = icon_count * icon_size + (icon_count - 1) * Inches(0.8)
        start_left = (SLIDE_W - total_w) / 2

        for idx in range(icon_count):
            icon_path = _find_icon(images_dir, slug, num, idx + 1)
            icon_left = start_left + idx * (icon_size + Inches(0.8))
            icon_top  = Inches(1.8)

            # Icon image or coloured box
            if icon_path.exists():
                slide.shapes.add_picture(str(icon_path), icon_left, icon_top, icon_size, icon_size)
            else:
                box = _rect(slide, icon_left, icon_top, icon_size, icon_size, C_ACCENT)

            # Label below icon (matching bullet)
            label = bullets[idx] if idx < len(bullets) else ""
            _textbox(slide, label,
                     icon_left - Inches(0.2), icon_top + icon_size + Inches(0.15),
                     icon_size + Inches(0.4), Inches(1.2),
                     size=14, color=C_LIGHT, align=PP_ALIGN.CENTER)

        # Remaining bullets as text below icons
        remaining = bullets[icon_count:]
        if remaining:
            box = slide.shapes.add_textbox(MARGIN, Inches(5.0), SLIDE_W - MARGIN * 2, Inches(2.0))
            tf  = box.text_frame
            tf.word_wrap = True
            for i, b in enumerate(remaining):
                p   = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                run = p.add_run()
                run.text = f"\u25cf  {b}"
                run.font.size  = Pt(15)
                run.font.color.rgb = C_LIGHT

    # ── Layout B: Nur Bullets → zentrierte Text-Liste ────────────────────────
    else:
        box = slide.shapes.add_textbox(MARGIN, Inches(1.5), SLIDE_W - MARGIN * 2, Inches(5.5))
        tf  = box.text_frame
        tf.word_wrap = True
        for i, b in enumerate(bullets):
            p   = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(12)
            run = p.add_run()
            run.text = f"\u25cf  {b}"
            run.font.size  = Pt(19)
            run.font.color.rgb = C_LIGHT

    _speaker_notes(slide, data.get("speaker_notes", ""))
    return slide


# ── Main ───────────────────────────────────────────────────────────────────────

def main():
    parser = argparse.ArgumentParser(description="Build PowerPoint from slide_structure.md")
    parser.add_argument("--version", type=int, default=1, help="Image iteration version (1 or 2)")
    parser.add_argument("--style", type=str, default="dark-professional",
                        choices=list(STYLES.keys()),
                        help="Visual style: dark-professional | light-modern | minimal | bold-creative")
    args = parser.parse_args()

    set_style(args.style)

    project_dir  = Path(__file__).parent.parent
    workspace    = project_dir / "workspace"
    images_dir   = project_dir / "output" / "images"
    out_dir      = project_dir / "output" / "presentations"
    out_dir.mkdir(parents=True, exist_ok=True)

    slug = _get_project_slug(workspace)

    try:
        data = parse_slide_structure(workspace)
    except FileNotFoundError as e:
        print(json.dumps({"ok": False, "error": str(e)}))
        sys.exit(1)
    except Exception as e:
        print(json.dumps({"ok": False, "error": f"Parse error: {e}"}))
        sys.exit(1)

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    slides = data.get("slides", [])
    if not slides:
        print(json.dumps({"ok": False, "error": "No slides found in slide_structure.md"}))
        sys.exit(1)

    for i, slide_data in enumerate(slides):
        num        = slide_data.get("number") or (i + 1)
        slide_type = slide_data.get("type", "content")
        has_image  = slide_data.get("has_image", True)
        img_path   = _find_image(images_dir, slug, num, slide_type)

        if slide_type == "title":
            build_title_slide(prs, slide_data, img_path if has_image else None)
        elif slide_type == "section":
            build_section_slide(prs, slide_data, img_path if has_image else None)
        elif slide_type == "closing":
            build_closing_slide(prs, slide_data, img_path if has_image else None)
        elif slide_type in ("data", "agenda", "quote"):
            build_data_slide(prs, slide_data, images_dir, slug)
        elif slide_type == "visual":
            build_visual_slide(prs, slide_data, img_path if has_image else None)
        else:
            build_content_slide(prs, slide_data, img_path if has_image else None)

    out_path = out_dir / f"{slug}_v{args.version}.pptx"
    prs.save(str(out_path))

    print(json.dumps({
        "ok":      True,
        "path":    str(out_path),
        "slides":  len(slides),
        "version": args.version,
    }, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
