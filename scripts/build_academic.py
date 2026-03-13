#!/usr/bin/env python3
"""
PresentationBanana — Academic Presentation Builder
====================================================
Custom builder for Baraghith paper with programmatic diagrams.
Uses atmospheric images only for title/section/closing slides.
Builds proper academic diagrams for content slides.
"""

import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print(json.dumps({"ok": False, "error": "pip install python-pptx"}))
    sys.exit(1)

# ── Dimensions ────────────────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Colors (dark-professional) ────────────────────────────────────────────────
C_BG      = RGBColor(0x12, 0x1A, 0x2E)
C_BG_MID  = RGBColor(0x1A, 0x27, 0x44)
C_BG_DARK = RGBColor(0x0D, 0x12, 0x20)
C_ACCENT  = RGBColor(0xF0, 0xAB, 0x00)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT   = RGBColor(0xC8, 0xD6, 0xE5)
C_MUTED   = RGBColor(0x7A, 0x8B, 0xA0)
C_BOX     = RGBColor(0x1E, 0x30, 0x50)
C_BOX2    = RGBColor(0x24, 0x3B, 0x5E)
C_RED     = RGBColor(0xE0, 0x4B, 0x4B)
C_GREEN   = RGBColor(0x3E, 0xB4, 0x89)
C_BLUE    = RGBColor(0x4A, 0x9E, 0xE0)
C_PURPLE  = RGBColor(0x9B, 0x6D, 0xD0)
C_TEAL    = RGBColor(0x2E, 0xA8, 0x9D)
C_ORANGE  = RGBColor(0xE8, 0x85, 0x3D)

LEVEL_COLORS = [C_ACCENT, C_BLUE, C_TEAL, C_PURPLE]  # L3, L2, L1, L0


# ── Helpers ───────────────────────────────────────────────────────────────────

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _bg(slide, color=C_BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color

def _rect(slide, left, top, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def _rounded_rect(slide, left, top, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.color.rgb = C_ACCENT
    s.line.width = Pt(1.5)
    return s

def _text(slide, text, left, top, w, h, size=18, bold=False, color=C_WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box

def _multitext(slide, lines, left, top, w, h, size=16, color=C_LIGHT, spacing=8):
    """Multiple paragraphs in one textbox."""
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(spacing)
        r = p.add_run()
        r.text = line
        r.font.size = Pt(size)
        r.font.color.rgb = color
    return box

def _labeled_box(slide, left, top, w, h, title, body, box_color=C_BOX, title_color=C_ACCENT, body_color=C_LIGHT):
    """A rounded box with title and body text."""
    s = _rounded_rect(slide, left, top, w, h, box_color)
    s.line.color.rgb = title_color
    s.line.width = Pt(1)
    _text(slide, title, left + Inches(0.15), top + Inches(0.08), w - Inches(0.3), Inches(0.35),
          size=13, bold=True, color=title_color)
    _text(slide, body, left + Inches(0.15), top + Inches(0.4), w - Inches(0.3), h - Inches(0.5),
          size=11, color=body_color)

def _arrow_down(slide, x, y1, y2, color=C_ACCENT):
    """Vertical arrow from (x,y1) to (x,y2)."""
    connector = slide.shapes.add_connector(1, x, y1, x, y2)  # 1=straight
    connector.line.color.rgb = color
    connector.line.width = Pt(2)

def _arrow_shape(slide, left, top, w, h, color=C_ACCENT):
    """Down arrow shape."""
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def _up_arrow_shape(slide, left, top, w, h, color=C_GREEN):
    """Up arrow shape."""
    s = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def _right_arrow(slide, left, top, w, h, color=C_ACCENT):
    """Right arrow shape."""
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def _chevron(slide, left, top, w, h, color=C_ACCENT):
    """Chevron/arrow shape."""
    s = slide.shapes.add_shape(MSO_SHAPE.CHEVRON, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def _notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text

def _img(slide, path, left, top, w, h):
    if path and Path(path).exists():
        slide.shapes.add_picture(str(path), left, top, w, h)

def _footer(slide, slide_num, total=19):
    """Add consistent footer with slide number and citation."""
    # Bottom accent line
    _rect(slide, Inches(0), Inches(7.25), SLIDE_W, Inches(0.04), C_ACCENT)
    # Slide number
    _text(slide, f"{slide_num} / {total}", Inches(12.0), Inches(7.28), Inches(1.2), Inches(0.22),
          size=9, color=C_MUTED, align=PP_ALIGN.RIGHT)
    # Citation
    _text(slide, "Baraghith (2025) · From Prompts to Populations",
          Inches(0.3), Inches(7.28), Inches(5), Inches(0.22),
          size=9, color=C_MUTED)

def _title_bar(slide, title, subtitle=None):
    """Standard title + accent line at top."""
    _text(slide, title, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
          size=28, bold=True, color=C_WHITE)
    _rect(slide, Inches(0.5), Inches(1.05), Inches(6.0), Inches(0.04), C_ACCENT)
    if subtitle:
        _text(slide, subtitle, Inches(0.5), Inches(1.15), Inches(12), Inches(0.5),
              size=14, color=C_MUTED)


# ── Image paths ───────────────────────────────────────────────────────────────

IMG_DIR = Path(__file__).parent.parent / "output" / "images"
SLUG = "from-prompts-to-populations"

def img(num, stype):
    return IMG_DIR / f"{SLUG}_s{num:02d}_{stype}.png"


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    """Title slide with atmospheric image."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _img(slide, img(1, "title"), Inches(7.4), Inches(0.3), Inches(5.6), Inches(6.9))
    _rect(slide, Inches(7.1), Inches(0.3), Inches(0.05), Inches(6.9), C_ACCENT)
    _text(slide, "From Prompts\nto Populations", Inches(0.5), Inches(1.2), Inches(6.4), Inches(2.5),
          size=40, bold=True)
    _rect(slide, Inches(0.5), Inches(3.8), Inches(4.0), Inches(0.06), C_ACCENT)
    _text(slide, "Cultural Evolution as a Meta-Frame\nfor AI Ontology and Individuation",
          Inches(0.5), Inches(4.1), Inches(6.4), Inches(1.5), size=20, color=C_LIGHT)
    _text(slide, "Baraghith (2025)", Inches(0.5), Inches(5.5), Inches(6), Inches(0.4),
          size=16, bold=True, color=C_ACCENT)
    _text(slide, "Presentation adapted from the original paper", Inches(0.5), Inches(5.9), Inches(6), Inches(0.4),
          size=12, color=C_MUTED)
    _notes(slide, "Based on Baraghith's paper proposing Cultural Evolutionary Theory as a unified framework for understanding AI individuation.")


def slide_02_central_question(prs):
    """Central question - clean text slide with key questions."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, "The Central Question")

    # Big quote
    _text(slide, '"How many kinds of AI are out there?"',
          Inches(1.0), Inches(1.8), Inches(11), Inches(1.0),
          size=30, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    # Key points as styled boxes
    points = [
        ("Engineers", "version by architecture + weights"),
        ("RL Researchers", "track behavioural policies"),
        ("Auditors", "need instance-level traceability"),
        ("Roboticists", "focus on embodiment + control loops"),
        ("Governance", "track purpose + organisational context"),
        ("Law", "requires accountable legal entities"),
    ]
    cols = 3
    bw, bh = Inches(3.6), Inches(1.1)
    gap_x, gap_y = Inches(0.4), Inches(0.3)
    start_x = Inches(0.7)
    start_y = Inches(3.2)

    for i, (who, what) in enumerate(points):
        col = i % cols
        row = i // cols
        x = start_x + col * (bw + gap_x)
        y = start_y + row * (bh + gap_y)
        _labeled_box(slide, x, y, bw, bh, who, what)

    _text(slide, "These criteria frequently conflict — and that is not a bug.",
          Inches(0.5), Inches(6.3), Inches(12), Inches(0.6),
          size=16, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    _notes(slide, "Baraghith's starting observation: different practitioners already individuate AI by different criteria, and these criteria frequently yield conflicting answers.")


def slide_03_roadmap(prs):
    """Paper structure overview."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, "Paper Structure")

    sections = [
        ("§2", "AI Ontology & Individuation", "Survey of philosophical positions\n+ six practical criteria + their conflicts", C_BLUE),
        ("§3", "CET as Meta-Frame", "Why Cultural Evolutionary Theory,\nnot just methodological pluralism?", C_TEAL),
        ("§4", "Four-Level Hierarchy", "Models → Systems → Episodes → Outputs\nwith V/T/S dynamics at each level", C_ACCENT),
        ("§5", "Mapping & Test Case", "Criteria mapped onto hierarchy;\nRLHF as multi-level test case", C_PURPLE),
        ("§6", "Conclusion & Governance", "One process, many levels;\nimplications for AI regulation", C_ORANGE),
    ]

    bw = Inches(2.2)
    bh = Inches(3.8)
    gap = Inches(0.2)
    start_x = Inches(0.5)
    start_y = Inches(1.7)

    for i, (sec, title, desc, color) in enumerate(sections):
        x = start_x + i * (bw + gap)
        s = _rounded_rect(slide, x, start_y, bw, bh, C_BOX)
        s.line.color.rgb = color
        s.line.width = Pt(2)
        # Section number circle
        _text(slide, sec, x + Inches(0.1), start_y + Inches(0.15), bw - Inches(0.2), Inches(0.4),
              size=22, bold=True, color=color, align=PP_ALIGN.CENTER)
        _rect(slide, x + Inches(0.3), start_y + Inches(0.6), bw - Inches(0.6), Inches(0.03), color)
        # Title
        _text(slide, title, x + Inches(0.12), start_y + Inches(0.75), bw - Inches(0.24), Inches(0.8),
              size=13, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
        # Description
        _text(slide, desc, x + Inches(0.12), start_y + Inches(1.6), bw - Inches(0.24), Inches(1.8),
              size=12, color=C_LIGHT, align=PP_ALIGN.CENTER)

        # Arrows between sections
        if i < 4:
            _right_arrow(slide, x + bw + Inches(0.02), start_y + Inches(1.5), Inches(0.16), Inches(0.3), color)

    _text(slide, "Core argument: the six individuation criteria are not rivals — they track entities at different levels of one evolutionary process.",
          Inches(0.5), Inches(6.0), Inches(12), Inches(0.8),
          size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    _footer(slide, 3)
    _notes(slide, "The paper proceeds from surveying existing positions, through introducing CET, to developing the four-level hierarchy and mapping everything onto it.")


def slide_04_ontology_section(prs):
    """Section divider with atmospheric image."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    p = img(3, "section")
    if p.exists():
        slide.shapes.add_picture(str(p), Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    _rect(slide, Inches(0), Inches(5.2), SLIDE_W, Inches(2.3), RGBColor(0, 0, 0))
    _rect(slide, Inches(0), Inches(5.2), SLIDE_W, Inches(0.04), C_ACCENT)
    _text(slide, "The Ontology Problem", Inches(0.8), Inches(5.35), Inches(11.7), Inches(1.0),
          size=34, bold=True, align=PP_ALIGN.CENTER)
    _text(slide, "Four philosophical positions on what AI \"is\"",
          Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.9),
          size=18, color=C_LIGHT, align=PP_ALIGN.CENTER)


def slide_05_four_views(prs):
    """Four philosophical positions as comparison boxes."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, "Four Philosophical Positions on AI Identity")

    views = [
        ("Ferrario (2025)", "Artifact Realism",
         "AI systems as genuine artifact kinds.\nIdentity via techno-function +\ntrustworthiness profile.",
         "Bounded individuals", C_BLUE),
        ("Cabitza et al. (2025)", "Cybork Concept",
         "Intelligence emerges from human-\nmachine-workflow ensembles.\nNot a property of bounded units.",
         "De-individuated networks", C_TEAL),
        ("Weinbaum & Veitas (2017)", "Process Ontology",
         "No pre-given agents — only processes\nof individuation. Boundaries, agents\nand goals co-emerge.",
         "Ongoing processes", C_PURPLE),
        ("Hawley (2019)", "Ontological Caution",
         "\"AI\" may not be a stable kind at all.\nCaution against reifying a moving\nconceptual frontier.",
         "Unstable category", C_ORANGE),
    ]

    bw = Inches(2.9)
    bh = Inches(3.6)
    gap = Inches(0.25)
    start_x = Inches(0.45)
    start_y = Inches(1.6)

    for i, (author, label, desc, stance, color) in enumerate(views):
        x = start_x + i * (bw + gap)
        # Box
        s = _rounded_rect(slide, x, start_y, bw, bh, C_BOX)
        s.line.color.rgb = color
        s.line.width = Pt(2)
        # Color bar at top
        _rect(slide, x + Inches(0.1), start_y + Inches(0.1), bw - Inches(0.2), Inches(0.06), color)
        # Author
        _text(slide, author, x + Inches(0.15), start_y + Inches(0.25), bw - Inches(0.3), Inches(0.35),
              size=11, color=C_MUTED)
        # Label
        _text(slide, label, x + Inches(0.15), start_y + Inches(0.55), bw - Inches(0.3), Inches(0.45),
              size=16, bold=True, color=color)
        # Description
        _text(slide, desc, x + Inches(0.15), start_y + Inches(1.1), bw - Inches(0.3), Inches(1.6),
              size=12, color=C_LIGHT)
        # Stance tag at bottom
        _text(slide, stance, x + Inches(0.15), start_y + bh - Inches(0.55), bw - Inches(0.3), Inches(0.4),
              size=11, bold=True, color=color, align=PP_ALIGN.CENTER)

    # Spectrum line showing individuation axis
    _rect(slide, Inches(0.5), Inches(5.6), Inches(12.3), Inches(0.04), C_MUTED)
    _text(slide, "← Strong individuation", Inches(0.5), Inches(5.7), Inches(3.0), Inches(0.3),
          size=11, color=C_BLUE)
    _text(slide, "Weak individuation →", Inches(9.8), Inches(5.7), Inches(3.0), Inches(0.3),
          size=11, color=C_ORANGE, align=PP_ALIGN.RIGHT)

    _text(slide, "Differ on: (1) How strongly they reify AI as individuals  (2) Where they locate agency",
          Inches(0.5), Inches(6.5), Inches(12), Inches(0.5),
          size=13, color=C_MUTED, align=PP_ALIGN.CENTER)
    _notes(slide, "These four positions differ along two dimensions: how strongly they reify AI as discrete individuals, and where they locate agency.")


def slide_06_six_criteria(prs):
    """Six individuation criteria as labeled diagram boxes."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, "Six Practical Individuation Criteria")

    criteria = [
        ("1. Architecture + Weights", "Model Instance Individuality",
         "Same file = same AI.\nNew weights (even tiny update) = new individual.\nAnalogy: two identical MP3 files are 'the same'.",
         C_BLUE),
        ("2. Policy Equivalence", "Behavioural Individuality",
         "Same behaviour = same AI.\nWhat it does, not how it's built.\nAnalogy: two musicians playing the same notes.",
         C_TEAL),
        ("3. Instance ID", "Session Individuality",
         "Same session/run = same AI.\nEach chat session is a unique individual.\nFine-grained accountability unit.",
         C_GREEN),
        ("4. Control Loop", "Embodiment Individuality",
         "Same body + environment = same AI.\nSwap the robot body = different agent.\nIdentity tracks the sensorimotor envelope.",
         C_PURPLE),
        ("5. Purpose + Organisation", "Functional Individuality",
         "Same function + context = same AI.\nBank A's credit scorer ≠ Bank B's,\neven if same base model.",
         C_ORANGE),
        ("6. Legal Status", "Normative Individuality",
         "Same accountable entity = same AI.\nThe entity that bears duties/liability.\nMost 'thick' and anthropomorphic criterion.",
         C_RED),
    ]

    cols = 3
    bw, bh = Inches(3.9), Inches(2.3)
    gap_x, gap_y = Inches(0.3), Inches(0.25)
    start_x = Inches(0.35)
    start_y = Inches(1.5)

    for i, (num_label, sub_label, desc, color) in enumerate(criteria):
        col = i % cols
        row = i // cols
        x = start_x + col * (bw + gap_x)
        y = start_y + row * (bh + gap_y)

        s = _rounded_rect(slide, x, y, bw, bh, C_BOX)
        s.line.color.rgb = color
        s.line.width = Pt(1.5)

        _text(slide, num_label, x + Inches(0.12), y + Inches(0.08), bw - Inches(0.24), Inches(0.3),
              size=13, bold=True, color=color)
        _text(slide, sub_label, x + Inches(0.12), y + Inches(0.35), bw - Inches(0.24), Inches(0.3),
              size=11, color=C_MUTED)
        _text(slide, desc, x + Inches(0.12), y + Inches(0.7), bw - Inches(0.24), Inches(1.5),
              size=11, color=C_LIGHT)

    _text(slide, "Each criterion is used in practice. They frequently yield conflicting answers.",
          Inches(0.5), Inches(6.7), Inches(12), Inches(0.5),
          size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    _notes(slide, "Engineers, lawyers, RL researchers, and roboticists all use different criteria — and they often disagree about what counts as 'the same AI'.")


def slide_07_conflicts(prs):
    """Criteria conflicts as specific conflict pairs."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, "When Criteria Conflict",
               "Same model ≠ same behaviour ≠ same legal entity")

    conflicts = [
        ("Model Instance", "vs", "Behavioural",
         "Same weights + different prompts\n→ same model, different behaviour",
         "Same test results + different weights\n→ different model, same behaviour",
         C_BLUE, C_TEAL),
        ("Model Instance", "vs", "Instance ID",
         "One model spawns millions of sessions\n→ one model, millions of 'individuals'",
         "New session, same weights\n→ new instance, same model",
         C_BLUE, C_GREEN),
        ("Control Loop", "vs", "Model Instance",
         "Same weights on different robot bodies\n→ same model, different agent",
         "Different weights, same body\n→ different model, same agent",
         C_PURPLE, C_BLUE),
        ("Instance ID", "vs", "Legal Personhood",
         "Millions of sessions map to\none accountable legal entity",
         "Same legal entity, different instances\n→ asymmetric nesting",
         C_GREEN, C_RED),
    ]

    bw = Inches(12.3)
    bh = Inches(1.05)
    start_x = Inches(0.5)
    start_y = Inches(1.7)
    gap = Inches(0.15)

    for i, (left, vs, right, ex1, ex2, lc, rc) in enumerate(conflicts):
        y = start_y + i * (bh + gap)
        # Background bar
        _rect(slide, start_x, y, bw, bh, C_BOX)
        # Left criterion
        _text(slide, left, start_x + Inches(0.2), y + Inches(0.05), Inches(2.0), Inches(0.35),
              size=13, bold=True, color=lc)
        # vs
        _text(slide, "vs", start_x + Inches(2.2), y + Inches(0.05), Inches(0.5), Inches(0.35),
              size=13, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
        # Right criterion
        _text(slide, right, start_x + Inches(2.7), y + Inches(0.05), Inches(2.0), Inches(0.35),
              size=13, bold=True, color=rc)
        # Examples
        _text(slide, ex1, start_x + Inches(0.2), y + Inches(0.4), Inches(5.5), Inches(0.6),
              size=11, color=C_LIGHT)
        _text(slide, ex2, start_x + Inches(6.2), y + Inches(0.4), Inches(5.5), Inches(0.6),
              size=11, color=C_LIGHT)

    _text(slide, '"These conflicts are not resolvable in general, because each criterion targets\na different explanatory job and a different level of organisation." — Baraghith',
          Inches(0.5), Inches(6.3), Inches(12), Inches(0.8),
          size=12, color=C_ACCENT, align=PP_ALIGN.CENTER)
    _notes(slide, "The six criteria cut along different axes and they often disagree. These conflicts cannot be resolved by choosing one 'true' criterion.")


def slide_08_cet_section(prs):
    """CET section divider with DNA→circuit image."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    p = img(7, "section")
    if p.exists():
        slide.shapes.add_picture(str(p), Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    _rect(slide, Inches(0), Inches(5.2), SLIDE_W, Inches(2.3), RGBColor(0, 0, 0))
    _rect(slide, Inches(0), Inches(5.2), SLIDE_W, Inches(0.04), C_ACCENT)
    _text(slide, "Cultural Evolutionary Theory", Inches(0.8), Inches(5.35), Inches(11.7), Inches(1.0),
          size=34, bold=True, align=PP_ALIGN.CENTER)
    _text(slide, "A unifying meta-frame, not just pluralism",
          Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.9),
          size=18, color=C_LIGHT, align=PP_ALIGN.CENTER)


def slide_09_why_cet(prs):
    """Why CET not pluralism - comparison diagram."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, "Why CET, Not Just Pluralism?")

    # Left side: Pluralism
    _rect(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.2), C_BOX)
    _text(slide, "Methodological Pluralism", Inches(0.8), Inches(1.6), Inches(5.2), Inches(0.5),
          size=20, bold=True, color=C_RED)
    _rect(slide, Inches(0.8), Inches(2.15), Inches(2.5), Inches(0.03), C_RED)

    plur_points = [
        '✓  "Use whichever criterion fits"',
        '✓  Descriptive tolerance of multiple views',
        '✗  Cannot explain WHY these criteria exist',
        '✗  Cannot explain WHY they conflict',
        '✗  Cannot trace causal connections',
        '✗  "Four separate lenses"',
    ]
    for i, pt in enumerate(plur_points):
        color = C_GREEN if pt.startswith('✓') else C_RED
        _text(slide, pt, Inches(0.8), Inches(2.5) + i * Inches(0.55), Inches(5.2), Inches(0.5),
              size=14, color=color)

    # Right side: CET
    _rect(slide, Inches(6.9), Inches(1.5), Inches(5.8), Inches(5.2), C_BOX)
    _text(slide, "Cultural Evolutionary Theory", Inches(7.2), Inches(1.6), Inches(5.2), Inches(0.5),
          size=20, bold=True, color=C_GREEN)
    _rect(slide, Inches(7.2), Inches(2.15), Inches(3.2), Inches(0.03), C_GREEN)

    cet_points = [
        '✓  Preserves insight: multiple criteria are legitimate',
        '✓  EXPLAINS why criteria exist (evolutionary levels)',
        '✓  PREDICTS which criteria cluster together',
        '✓  Traces upward + downward causation',
        '✓  Maps interventions to cascading effects',
        '✓  "One coupled dynamical system"',
    ]
    for i, pt in enumerate(cet_points):
        _text(slide, pt, Inches(7.2), Inches(2.5) + i * Inches(0.55), Inches(5.2), Inches(0.5),
              size=14, color=C_GREEN)

    # Divider
    _rect(slide, Inches(6.45), Inches(1.7), Inches(0.04), Inches(4.8), C_ACCENT)
    _text(slide, "→", Inches(6.15), Inches(3.6), Inches(0.6), Inches(0.6),
          size=30, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    _notes(slide, "CET does not reject pluralism but subsumes it: it preserves the insight that multiple criteria are legitimate while adding a structural explanation of why they exist, how they relate, and what happens when we intervene at different levels.")


def slide_10_hierarchy_overview(prs):
    """THE core diagram: four-level hierarchy."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, "The Four-Level Evolutionary Hierarchy",
               "Each level has its own variation, transmission, and selection dynamics")

    levels = [
        ("Level 3", "Model Lineages", "GPT-4, LLaMA families, fine-tunes, merges",
         "Variation: new architectures, hyperparameters\nTransmission: cloning, forking, distillation\nSelection: benchmarks, safety evals, adoption",
         "≈ Genetic lineages", C_ACCENT),
        ("Level 2", "Configured Systems", "Model + prompts + tools + guardrails + oversight",
         "Variation: swapping models, changing tools\nTransmission: templating, copying configs\nSelection: performance, cost, regulation",
         "≈ Organisms", C_BLUE),
        ("Level 1", "Episodes & Instances", "Individual chat sessions, API calls, robot runs",
         "Variation: different prompts, contexts, seeds\nTransmission: patterns → training data\nSelection: user feedback, ratings, monitoring",
         "≈ Interactions", C_TEAL),
        ("Level 0", "Circulating Outputs", "Texts, images, code, scores, decisions",
         "Variation: stochastic sampling, editing\nTransmission: sharing, quoting, scraping\nSelection: attention, platform ranking",
         "≈ Cultural memes", C_PURPLE),
    ]

    bw = Inches(12.3)
    bh = Inches(1.15)
    start_x = Inches(0.5)
    start_y = Inches(1.8)
    gap = Inches(0.12)

    for i, (lvl, name, examples, dynamics, analogy, color) in enumerate(levels):
        y = start_y + i * (bh + gap)
        # Background
        _rect(slide, start_x, y, bw, bh, C_BOX)
        # Color bar left
        _rect(slide, start_x, y, Inches(0.08), bh, color)
        # Level label
        _text(slide, lvl, start_x + Inches(0.2), y + Inches(0.05), Inches(1.0), Inches(0.35),
              size=14, bold=True, color=color)
        # Name
        _text(slide, name, start_x + Inches(1.3), y + Inches(0.05), Inches(2.5), Inches(0.35),
              size=14, bold=True, color=C_WHITE)
        # Examples
        _text(slide, examples, start_x + Inches(1.3), y + Inches(0.38), Inches(3.2), Inches(0.7),
              size=11, color=C_MUTED)
        # Dynamics
        _text(slide, dynamics, start_x + Inches(4.8), y + Inches(0.05), Inches(5.0), Inches(1.0),
              size=11, color=C_LIGHT)
        # Bio analogy
        _text(slide, analogy, start_x + Inches(10.2), y + Inches(0.05), Inches(1.8), Inches(0.35),
              size=11, bold=True, color=color, align=PP_ALIGN.RIGHT)

        # Down arrow between levels
        if i < 3:
            ay = y + bh
            _arrow_shape(slide, Inches(1.2), ay, Inches(0.3), gap, color)

    _notes(slide, "The four levels correspond roughly to biological evolution: Level 3 = genetic lineages, Level 2 = organisms, Level 1 = organism-environment interactions, Level 0 = behavioural traces.")


def slide_11_level3(prs):
    """Level 3: Models and Lineages - tree diagram."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, "Level 3: Model Lineages")

    # Show a model family tree
    _text(slide, "Model lineages accumulate changes and branch — like species lineages",
          Inches(0.5), Inches(1.3), Inches(12), Inches(0.4),
          size=14, color=C_MUTED)

    # Tree structure using boxes and connectors
    # Root
    _labeled_box(slide, Inches(5.0), Inches(1.9), Inches(3.0), Inches(0.7),
                 "Base Model (pretrained)", "Architecture + weights θ", C_BOX, C_ACCENT)

    # Level 1 branches
    branches = [
        (Inches(0.8), "Fine-tune A\n(medical data)"),
        (Inches(4.0), "Fine-tune B\n(code generation)"),
        (Inches(7.2), "Distillation C\n(smaller, faster)"),
        (Inches(10.3), "LoRA Adapter D\n(task-specific)"),
    ]

    # Root center X
    root_cx = Inches(5.0) + Inches(1.5)  # 6.5"
    root_bottom = Inches(2.6)
    branch_top = Inches(3.3)

    # Vertical trunk from root down to branch level
    _rect(slide, root_cx - Inches(0.02), root_bottom, Inches(0.04), Inches(0.35), C_MUTED)
    # Horizontal bar spanning all branches
    first_cx = Inches(0.8) + Inches(1.25)  # 2.05"
    last_cx = Inches(10.3) + Inches(1.25)  # 11.55"
    bar_y = root_bottom + Inches(0.33)
    _rect(slide, first_cx, bar_y, last_cx - first_cx, Inches(0.04), C_MUTED)

    for x, label in branches:
        _labeled_box(slide, x, branch_top, Inches(2.5), Inches(0.8), label, "", C_BOX2, C_BLUE)
        # Vertical drop from horizontal bar to branch top
        branch_cx = x + Inches(1.25)
        _rect(slide, branch_cx - Inches(0.02), bar_y, Inches(0.04), branch_top - bar_y, C_MUTED)

    # Second generation with parent connections
    sub = [
        (Inches(0.3), "Merged A+B", Inches(0.8)),     # parent: Fine-tune A
        (Inches(2.6), "Quantised A", Inches(0.8)),     # parent: Fine-tune A
        (Inches(7.0), "C deployed\nin Hospital", Inches(7.2)),  # parent: Distillation C
        (Inches(9.5), "D v2\n(retrained)", Inches(10.3)),       # parent: LoRA D
    ]
    branch_bottom = Inches(4.1)
    sub_top = Inches(4.8)

    for x, label, parent_x in sub:
        _labeled_box(slide, x, sub_top, Inches(2.2), Inches(0.7), label, "", C_BOX2, C_TEAL)
        # Connect from parent branch center down to sub
        parent_cx = parent_x + Inches(1.25)
        sub_cx = x + Inches(1.1)
        # Vertical from parent bottom
        _rect(slide, parent_cx - Inches(0.02), branch_bottom, Inches(0.04), Inches(0.35), C_MUTED)
        # Horizontal to sub center if different
        conn_y = branch_bottom + Inches(0.33)
        if abs((parent_cx - sub_cx) / 914400) > 0.3:
            min_x = min(parent_cx, sub_cx)
            max_x = max(parent_cx, sub_cx)
            _rect(slide, min_x, conn_y, max_x - min_x, Inches(0.04), C_MUTED)
        # Vertical drop to sub
        _rect(slide, sub_cx - Inches(0.02), conn_y, Inches(0.04), sub_top - conn_y, C_MUTED)

    # Key points
    bullets = [
        "● Variation: new architectures, hyperparameters, training data, fine-tuning",
        "● Transmission: model copying, forking, API-mediated deployment",
        "● Selection: benchmarks, safety evaluations, adoption, licensing decisions",
        '● "What a model \'is\' — functionally — is its parameterization" (Baraghith)',
    ]
    _multitext(slide, bullets, Inches(0.5), Inches(5.8), Inches(12), Inches(1.5),
               size=13, color=C_LIGHT)
    _notes(slide, "Level 3 is where model-instance individuality and behavioural individuality most naturally belong. Examples: GPT-4.1, LLaMA-3-8B finetuned on biomedical abstracts.")


def slide_12_level2(prs):
    """Level 2: Configured Systems - architecture diagram."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, "Level 2: Configured Systems")

    _text(slide, "The socio-technical 'bundle' that acts in the world — like organisms in an environment",
          Inches(0.5), Inches(1.3), Inches(12), Inches(0.4), size=14, color=C_MUTED)

    # Central system diagram
    # Outer boundary
    s = _rounded_rect(slide, Inches(2.5), Inches(2.0), Inches(8.0), Inches(3.5), C_BOX)
    s.line.color.rgb = C_BLUE
    s.line.width = Pt(2)
    _text(slide, "Configured System (Level 2 Individual)", Inches(2.7), Inches(2.05), Inches(7.5), Inches(0.4),
          size=12, bold=True, color=C_BLUE)

    # Inner components
    components = [
        (Inches(2.8), Inches(2.6), Inches(2.3), Inches(1.0), "Model Weights\n(from Level 3)", C_ACCENT),
        (Inches(5.3), Inches(2.6), Inches(2.3), Inches(1.0), "System Prompt\n+ Safety Rules", C_TEAL),
        (Inches(7.8), Inches(2.6), Inches(2.3), Inches(1.0), "Tools + APIs\n(web, code, DB)", C_GREEN),
        (Inches(2.8), Inches(3.8), Inches(2.3), Inches(1.0), "Guardrails\n+ Filters", C_RED),
        (Inches(5.3), Inches(3.8), Inches(2.3), Inches(1.0), "Data Pipelines\n+ RAG Index", C_PURPLE),
        (Inches(7.8), Inches(3.8), Inches(2.3), Inches(1.0), "Human Oversight\n+ Escalation", C_ORANGE),
    ]

    for x, y, w, h, label, color in components:
        box = _rounded_rect(slide, x, y, w, h, C_BOX2)
        box.line.color.rgb = color
        _text(slide, label, x + Inches(0.1), y + Inches(0.15), w - Inches(0.2), h - Inches(0.3),
              size=11, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # Examples
    _text(slide, 'Examples: "The sepsis-alert system v3 in Hospital A"  •  "TikTok\'s 2025 recommendation stack"  •  "Bank B\'s credit scoring deployment"',
          Inches(0.5), Inches(5.7), Inches(12), Inches(0.5),
          size=12, color=C_MUTED, align=PP_ALIGN.CENTER)

    bullets = [
        "● Selected by: performance, cost, usability, regulatory compliance",
        "● This is where Purpose and Control-loop individuality primarily live",
    ]
    _multitext(slide, bullets, Inches(0.5), Inches(6.2), Inches(12), Inches(1.0),
               size=13, color=C_LIGHT)
    _notes(slide, "Level 2 systems are the comparatively stable, bounded configurations that succeed or fail — like organisms in biological evolution. They are socio-technical ensembles.")


def slide_13_level1(prs):
    """Level 1: Episodes - with chat diagram."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, "Level 1: Episodes & Instances")

    _text(slide, "Short-lived vehicles through which selection signals and training data are generated",
          Inches(0.5), Inches(1.3), Inches(12), Inches(0.4), size=14, color=C_MUTED)

    # Episode flow diagram
    # System box
    s = _rounded_rect(slide, Inches(4.5), Inches(2.0), Inches(4.5), Inches(1.2), C_BOX)
    s.line.color.rgb = C_TEAL
    _text(slide, "Configured System (Level 2)", Inches(4.7), Inches(2.15), Inches(4.0), Inches(0.4),
          size=12, bold=True, color=C_TEAL, align=PP_ALIGN.CENTER)
    _text(slide, "model + prompts + tools + guardrails",
          Inches(4.7), Inches(2.5), Inches(4.0), Inches(0.4),
          size=11, color=C_MUTED, align=PP_ALIGN.CENTER)

    # Episodes
    episodes = [
        (Inches(1.0), "Episode A\nUser asks about\nweather → response"),
        (Inches(4.5), "Episode B\nDeveloper tests\ncode generation"),
        (Inches(8.0), "Episode C\nUser complains\n→ escalation"),
    ]
    for x, label in episodes:
        s = _rounded_rect(slide, x, Inches(3.8), Inches(3.0), Inches(1.2), C_BOX2)
        s.line.color.rgb = C_GREEN
        _text(slide, label, x + Inches(0.15), Inches(3.9), Inches(2.7), Inches(1.0),
              size=11, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # Feedback arrows going up
    _text(slide, "positive rating", Inches(1.8), Inches(5.1), Inches(1.5), Inches(0.3),
          size=11, color=C_GREEN, align=PP_ALIGN.CENTER)
    _text(slide, "continued use", Inches(5.3), Inches(5.1), Inches(1.5), Inches(0.3),
          size=11, color=C_GREEN, align=PP_ALIGN.CENTER)
    _text(slide, "complaint flag", Inches(8.8), Inches(5.1), Inches(1.5), Inches(0.3),
          size=11, color=C_RED, align=PP_ALIGN.CENTER)

    _text(slide, "Selection signals", Inches(4.5), Inches(5.5), Inches(4.5), Inches(0.4),
          size=16, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    bullets = [
        "● Each episode = unique context, prompt, random seed → different trajectory",
        "● Where Instance ID individuality naturally belongs",
        "● Episodes don't reproduce — but patterns feed into training data for Level 3",
    ]
    _multitext(slide, bullets, Inches(0.5), Inches(6.0), Inches(12), Inches(1.2),
               size=13, color=C_LIGHT)
    _notes(slide, "Level 1 episodes sit at the interface between micro-level interaction and macro-level evolution: where variation and selection are most immediately experienced by users.")


def slide_14_level0(prs):
    """Level 0: Outputs - with image."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, "Level 0: Circulating Outputs")

    # Use atmospheric image on right
    _img(slide, img(13, "content"), Inches(8.0), Inches(0.0), Inches(5.3), SLIDE_H)
    _rect(slide, Inches(7.94), Inches(0.0), Inches(0.06), SLIDE_H, C_ACCENT)

    bullets = [
        '● Outputs = texts, images, code, scores,\n   decisions produced by AI systems',
        '● "The cultural memes of AI" — discrete\n   transmissible items (Boyd & Richerson 1985)',
        '● Not AI individuals themselves — but where\n   cultural selection is most directly visible',
        '● Variation: stochastic sampling, different\n   prompts, human modifications',
        '● Reproduction: copying, reposting, scraping\n   into training data',
        '● Selection: attention, sharing, platform\n   ranking, moderation',
        '● KEY: Where AI-driven cultural evolution\n   merges into ordinary human cultural evolution',
    ]
    _multitext(slide, bullets, Inches(0.5), Inches(1.5), Inches(7.2), Inches(5.5),
               size=14, color=C_LIGHT, spacing=10)
    _notes(slide, "Level 0 outputs are the primary carriers of provenance and influence across levels, traceable back to particular lineages and deployments, and forward into new training sets.")


def slide_15_cross_level(prs):
    """Cross-level causal pathways - THE key diagram."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, "Cross-Level Causal Pathways",
               "Downward causation constrains — Upward causation aggregates and reshapes")

    # Four level boxes
    levels = [
        ("Level 3: Models", "months–years", C_ACCENT),
        ("Level 2: Systems", "days–weeks", C_BLUE),
        ("Level 1: Episodes", "real-time", C_TEAL),
        ("Level 0: Outputs", "seconds–years", C_PURPLE),
    ]

    box_w = Inches(3.5)
    box_h = Inches(0.9)
    center_x = Inches(4.9)
    start_y = Inches(1.8)
    gap = Inches(0.5)

    for i, (name, timescale, color) in enumerate(levels):
        y = start_y + i * (box_h + gap)
        s = _rounded_rect(slide, center_x, y, box_w, box_h, C_BOX)
        s.line.color.rgb = color
        s.line.width = Pt(2)
        _text(slide, name, center_x + Inches(0.15), y + Inches(0.08), box_w - Inches(0.3), Inches(0.4),
              size=15, bold=True, color=color)
        _text(slide, timescale, center_x + Inches(0.15), y + Inches(0.45), box_w - Inches(0.3), Inches(0.3),
              size=11, color=C_MUTED)

    # Downward causation (left side) — thick vertical line + arrows
    down_x = Inches(3.8)
    _rect(slide, down_x, start_y + box_h, Inches(0.06), (box_h + gap) * 3 - gap, C_RED)
    _text(slide, "CONSTRAINS ↓", Inches(2.0), start_y + Inches(0.2), Inches(1.7), Inches(0.4),
          size=13, bold=True, color=C_RED, align=PP_ALIGN.RIGHT)
    for i in range(3):
        y1 = start_y + (i + 1) * (box_h + gap) + Inches(0.2)
        _arrow_shape(slide, down_x - Inches(0.1), y1, Inches(0.25), Inches(0.35), C_RED)

    # Upward aggregation (right side) — thick vertical line + arrows
    up_x = Inches(9.0)
    _rect(slide, up_x, start_y + box_h, Inches(0.06), (box_h + gap) * 3 - gap, C_GREEN)
    _text(slide, "↑ AGGREGATES", Inches(9.3), start_y + Inches(0.2), Inches(2.0), Inches(0.4),
          size=13, bold=True, color=C_GREEN)
    for i in range(3):
        y1 = start_y + i * (box_h + gap) + box_h + Inches(0.05)
        s = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, up_x - Inches(0.1), y1, Inches(0.25), Inches(0.35))
        s.fill.solid()
        s.fill.fore_color.rgb = C_GREEN
        s.line.fill.background()

    # Bypass arrow L0 → L3 — prominent curved path on far left
    bypass_x = Inches(1.5)
    # Vertical bar representing bypass
    bottom_y = start_y + 3 * (box_h + gap) + box_h / 2
    top_y = start_y + box_h / 2
    _rect(slide, bypass_x, top_y, Inches(0.08), bottom_y - top_y, C_ACCENT)
    # Horizontal connectors to boxes
    _rect(slide, bypass_x, bottom_y - Inches(0.02), center_x - bypass_x, Inches(0.04), C_ACCENT)
    _rect(slide, bypass_x, top_y - Inches(0.02), center_x - bypass_x, Inches(0.04), C_ACCENT)
    # Arrow head at top
    s = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, bypass_x - Inches(0.1), top_y - Inches(0.4), Inches(0.3), Inches(0.4))
    s.fill.solid()
    s.fill.fore_color.rgb = C_ACCENT
    s.line.fill.background()
    # Label
    _text(slide, "DIRECT\nFEEDBACK\nL0 → L3", Inches(0.1), Inches(3.3), Inches(1.3), Inches(1.2),
          size=12, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    _text(slide, "Outputs re-enter\ntraining corpora", Inches(0.1), Inches(4.5), Inches(1.3), Inches(0.6),
          size=11, color=C_LIGHT, align=PP_ALIGN.CENTER)

    # Timescale note
    _text(slide, "Temporal asymmetry: fast-changing lower levels can outrun slow-changing higher levels",
          Inches(2.0), Inches(6.5), Inches(10), Inches(0.5),
          size=13, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    _notes(slide, "Three key features of cross-level causal pathways: (1) asymmetric, (2) non-adjacent connections matter (L0→L3 bypass), (3) different timescales create temporal disequilibrium.")


def slide_16_rlhf(prs):
    """RLHF feedback loop - cycle diagram."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, 'Test Case: The RLHF Feedback Loop',
               '"No single level owns the evolutionary process"')

    # Four nodes in a diamond/cycle
    # Top: L3 Models
    nodes = [
        (Inches(5.2), Inches(1.8), "L3: Models", "Accumulated data\nreshapes weights via\npolicy optimisation", C_ACCENT),
        (Inches(9.0), Inches(3.5), "L2: Systems", "Config filters which\nepisodes occur;\nmediates feedback", C_BLUE),
        (Inches(5.2), Inches(5.0), "L1: Episodes", "Generate varied queries,\nresponses, and\nselection signals", C_TEAL),
        (Inches(1.2), Inches(3.5), "L0: Outputs", "Circulate beyond system;\nenter competitors'\ntraining corpora", C_PURPLE),
    ]

    nw, nh = Inches(2.8), Inches(1.3)
    for x, y, label, desc, color in nodes:
        s = _rounded_rect(slide, x, y, nw, nh, C_BOX)
        s.line.color.rgb = color
        s.line.width = Pt(2)
        _text(slide, label, x + Inches(0.1), y + Inches(0.05), nw - Inches(0.2), Inches(0.35),
              size=14, bold=True, color=color)
        _text(slide, desc, x + Inches(0.1), y + Inches(0.4), nw - Inches(0.2), Inches(0.8),
              size=11, color=C_LIGHT)

    # Connecting lines between nodes (clockwise: L3→L2, L2→L1, L1→L0, L0→L3)
    # L3 → L2 (top-right to right): horizontal connector
    _rect(slide, Inches(8.0), Inches(2.45), Inches(1.0), Inches(0.04), C_ACCENT)
    _rect(slide, Inches(8.96), Inches(2.45), Inches(0.04), Inches(1.1), C_ACCENT)
    _text(slide, "constrains", Inches(8.0), Inches(2.0), Inches(1.0), Inches(0.3),
          size=11, bold=True, color=C_MUTED, align=PP_ALIGN.CENTER)

    # L2 → L1 (right to bottom): vertical connector
    _rect(slide, Inches(9.36), Inches(4.8), Inches(0.04), Inches(0.7), C_ACCENT)
    _rect(slide, Inches(8.0), Inches(5.46), Inches(1.4), Inches(0.04), C_ACCENT)
    _text(slide, "feedback", Inches(9.4), Inches(4.9), Inches(1.5), Inches(0.3),
          size=11, bold=True, color=C_MUTED)

    # L1 → L0 (bottom to left): horizontal connector
    _rect(slide, Inches(4.0), Inches(5.66), Inches(1.2), Inches(0.04), C_ACCENT)
    _rect(slide, Inches(4.0), Inches(4.8), Inches(0.04), Inches(0.9), C_ACCENT)
    _text(slide, "generates", Inches(4.0), Inches(5.8), Inches(1.2), Inches(0.3),
          size=11, bold=True, color=C_MUTED, align=PP_ALIGN.CENTER)

    # L0 → L3 (left to top): vertical connector
    _rect(slide, Inches(2.56), Inches(2.1), Inches(0.04), Inches(1.4), C_ACCENT)
    _rect(slide, Inches(2.56), Inches(2.1), Inches(2.7), Inches(0.04), C_ACCENT)
    _text(slide, "re-enters", Inches(1.3), Inches(2.6), Inches(1.5), Inches(0.3),
          size=11, bold=True, color=C_MUTED, align=PP_ALIGN.CENTER)

    _text(slide, "All four levels contribute simultaneously — this is the point, not a defect.",
          Inches(0.5), Inches(6.7), Inches(12), Inches(0.5),
          size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    _notes(slide, "The RLHF process demonstrates how no single level owns the evolutionary process. Altering the feedback interface at Level 1 changes selection resolution. Changing guardrails at Level 2 changes the selection pool.")


def slide_17_mapping(prs):
    """Mapping criteria onto hierarchy - TABLE."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, "Mapping Individuation Criteria onto the Hierarchy",
               "The six criteria are not rivals — they are partial perspectives on different levels")

    # Table data
    rows = [
        ("Level 3\nModels", "Model-instance\nindividuality", "Behavioural/policy\nindividuality", "Where heritable information\nis stored and copied", C_ACCENT),
        ("Level 2\nSystems", "Purpose\nindividuality", "Control-loop\nindividuality", "Where functional integration\nproduces bounded units", C_BLUE),
        ("Level 1\nEpisodes", "Instance ID\nindividuality", "—", "Where selection events\nand training data originate", C_TEAL),
        ("Overlay", "Legal personhood\nindividuality", "—", "Normative overlay on\nany level (mainly L2)", C_PURPLE),
    ]

    # Header row
    headers = ["Level", "Primary Criterion", "Secondary Criterion", "Evolutionary Role"]
    hw = [Inches(1.8), Inches(2.8), Inches(2.8), Inches(4.2)]
    hx = [Inches(0.5)]
    for i in range(1, len(hw)):
        hx.append(hx[-1] + hw[i-1] + Inches(0.1))

    header_y = Inches(1.8)
    rh = Inches(0.5)

    for i, (h, w, x) in enumerate(zip(headers, hw, hx)):
        _rect(slide, x, header_y, w, rh, C_ACCENT)
        _text(slide, h, x + Inches(0.1), header_y + Inches(0.08), w - Inches(0.2), rh - Inches(0.15),
              size=13, bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER)

    # Data rows
    row_h = Inches(1.05)
    row_gap = Inches(0.08)
    start_y = header_y + rh + Inches(0.1)

    for ri, (level, crit1, crit2, role, color) in enumerate(rows):
        y = start_y + ri * (row_h + row_gap)
        for ci, (w, x) in enumerate(zip(hw, hx)):
            _rect(slide, x, y, w, row_h, C_BOX)
        # Color bar
        _rect(slide, hx[0], y, Inches(0.06), row_h, color)
        # Level
        _text(slide, level, hx[0] + Inches(0.15), y + Inches(0.1), hw[0] - Inches(0.3), row_h - Inches(0.2),
              size=13, bold=True, color=color)
        # Criteria
        _text(slide, crit1, hx[1] + Inches(0.15), y + Inches(0.1), hw[1] - Inches(0.3), row_h - Inches(0.2),
              size=12, color=C_WHITE)
        if crit2:
            _text(slide, crit2, hx[2] + Inches(0.15), y + Inches(0.1), hw[2] - Inches(0.3), row_h - Inches(0.2),
                  size=12, color=C_WHITE)
        # Role
        _text(slide, role, hx[3] + Inches(0.15), y + Inches(0.1), hw[3] - Inches(0.3), row_h - Inches(0.2),
              size=12, color=C_LIGHT)

    _text(slide, '"The evolutionary hierarchy predicts this clustering: individuation criteria group around levels\nwhere variation, transmission, or selection are most salient." — Baraghith',
          Inches(0.5), Inches(6.6), Inches(12), Inches(0.7),
          size=12, color=C_ACCENT, align=PP_ALIGN.CENTER)
    _notes(slide, "This mapping resolves the apparent conflicts: the criteria do not compete because they target different evolutionary levels.")


def slide_18_governance(prs):
    """Governance implications with cascade diagram."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, "Implications for AI Governance",
               "Different interventions target different levels with predictable cascading effects")

    interventions = [
        ("Regulate Training\n(Level 3)", "Restrict training data,\ncontrol open-weight release",
         "Cuts off variation at the source;\nconstrains entire downstream\nevolutionary tree",
         "broadest impact", C_ACCENT),
        ("Regulate Systems\n(Level 2)", "Require human oversight,\ncertification, auditing",
         "Constrains viable deployment\ndesigns; alters which selection\nsignals reach training",
         "targeted impact", C_BLUE),
        ("Regulate Episodes\n(Level 1)", "Feedback granularity,\nuser consent mechanisms",
         "Changes resolution of\nselection signals; affects\npreference data quality",
         "fine-grained", C_TEAL),
        ("Regulate Outputs\n(Level 0)", "Mandatory watermarks,\ncontent moderation",
         "Changes selective environment;\nreshapes which system configs\nremain competitive",
         "environmental", C_PURPLE),
    ]

    bw = Inches(2.65)
    bh = Inches(3.8)
    gap = Inches(0.5)
    start_x = Inches(0.5)
    start_y = Inches(1.7)

    for i, (label, what, effect, scope, color) in enumerate(interventions):
        x = start_x + i * (bw + gap)
        # Box
        s = _rounded_rect(slide, x, start_y, bw, bh, C_BOX)
        s.line.color.rgb = color
        s.line.width = Pt(2)
        # Color top bar
        _rect(slide, x + Inches(0.1), start_y + Inches(0.1), bw - Inches(0.2), Inches(0.05), color)
        # Label
        _text(slide, label, x + Inches(0.15), start_y + Inches(0.2), bw - Inches(0.3), Inches(0.6),
              size=14, bold=True, color=color)
        # What
        _text(slide, what, x + Inches(0.15), start_y + Inches(0.85), bw - Inches(0.3), Inches(0.7),
              size=11, color=C_MUTED)
        # Effect
        _text(slide, effect, x + Inches(0.15), start_y + Inches(1.7), bw - Inches(0.3), Inches(1.2),
              size=11, color=C_LIGHT)
        # Scope tag
        _text(slide, scope, x + Inches(0.15), start_y + bh - Inches(0.55), bw - Inches(0.3), Inches(0.4),
              size=11, bold=True, color=color, align=PP_ALIGN.CENTER)

    # Cascade arrows — properly centered between columns
    for i in range(3):
        arrow_x = start_x + (i + 1) * bw + i * gap + Inches(0.08)
        arrow_y = start_y + bh / 2 - Inches(0.15)
        s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, arrow_x, arrow_y, Inches(0.35), Inches(0.3))
        s.fill.solid()
        s.fill.fore_color.rgb = C_ACCENT
        s.line.fill.background()

    _text(slide, "CET provides the vocabulary for tracking these cascades — unstructured pluralism does not.",
          Inches(0.5), Inches(6.3), Inches(12), Inches(0.5),
          size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    _notes(slide, "A key practical payoff: different governance interventions target different levels with predictable cascading effects. CET provides the vocabulary for tracking these cascades.")


def slide_19_closing(prs):
    """Closing slide with atmospheric image."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _img(slide, img(18, "closing"), Inches(0), Inches(0), Inches(5.8), SLIDE_H)
    _rect(slide, Inches(5.8), Inches(0.3), Inches(0.05), Inches(6.9), C_ACCENT)

    _text(slide, "One Process,\nMany Levels", Inches(6.3), Inches(1.5), Inches(6.5), Inches(2.5),
          size=36, bold=True)
    _rect(slide, Inches(6.3), Inches(3.8), Inches(3.0), Inches(0.06), C_ACCENT)
    _text(slide, "AI individuation is not a choice between\ncriteria — it's recognising where each\nbelongs in the evolutionary hierarchy.",
          Inches(6.3), Inches(4.1), Inches(6.5), Inches(1.5), size=18, color=C_LIGHT)

    _text(slide, "Framework is methodologically independent\nof debates about AI consciousness or sentience.",
          Inches(6.3), Inches(5.5), Inches(6.5), Inches(1.0), size=13, color=C_MUTED)

    _footer(slide, 19)
    _notes(slide, "Conclusion: Rather than selecting one 'true' criterion, CET provides a meta-frame where all six criteria are partial perspectives on a single multi-level evolutionary process.")


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    project_dir = Path(__file__).parent.parent
    out_dir = project_dir / "output" / "presentations"
    out_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    TOTAL = 19

    # Build all slides (19 total with new roadmap slide)
    slide_01_title(prs)

    slide_02_central_question(prs)
    _footer(prs.slides[-1], 2, TOTAL)

    slide_03_roadmap(prs)  # has internal _footer call

    slide_04_ontology_section(prs)
    _footer(prs.slides[-1], 4, TOTAL)

    slide_05_four_views(prs)
    _footer(prs.slides[-1], 5, TOTAL)

    slide_06_six_criteria(prs)
    _footer(prs.slides[-1], 6, TOTAL)

    slide_07_conflicts(prs)
    _footer(prs.slides[-1], 7, TOTAL)

    slide_08_cet_section(prs)
    _footer(prs.slides[-1], 8, TOTAL)

    slide_09_why_cet(prs)
    _footer(prs.slides[-1], 9, TOTAL)

    slide_10_hierarchy_overview(prs)
    _footer(prs.slides[-1], 10, TOTAL)

    slide_11_level3(prs)
    _footer(prs.slides[-1], 11, TOTAL)

    slide_12_level2(prs)
    _footer(prs.slides[-1], 12, TOTAL)

    slide_13_level1(prs)
    _footer(prs.slides[-1], 13, TOTAL)

    slide_14_level0(prs)
    _footer(prs.slides[-1], 14, TOTAL)

    slide_15_cross_level(prs)
    _footer(prs.slides[-1], 15, TOTAL)

    slide_16_rlhf(prs)
    _footer(prs.slides[-1], 16, TOTAL)

    slide_17_mapping(prs)
    _footer(prs.slides[-1], 17, TOTAL)

    slide_18_governance(prs)
    _footer(prs.slides[-1], 18, TOTAL)

    slide_19_closing(prs)

    out_path = out_dir / "from-prompts-to-populations_v3.pptx"
    prs.save(str(out_path))

    print(json.dumps({
        "ok": True,
        "path": str(out_path),
        "slides": 19,
        "version": 3,
    }, indent=2))


if __name__ == "__main__":
    main()
