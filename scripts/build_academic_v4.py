#!/usr/bin/env python3
"""
PresentationBanana — Academic Presentation Builder v4
======================================================
Strict academic rules: 1 idea/slide, max 20 words, declarative titles,
diagrams over text, Imagen only for title/section/closing.
Bull & Bear critique-informed design.
"""

import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.util import Inches, Pt, Emu
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print(json.dumps({"ok": False, "error": "pip install python-pptx"}))
    sys.exit(1)

# ── Dimensions ────────────────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL = 19

# ── Colors (dark-professional) ────────────────────────────────────────────────
C_BG      = RGBColor(0x12, 0x1A, 0x2E)
C_BG_MID  = RGBColor(0x1A, 0x27, 0x44)
C_BG_DARK = RGBColor(0x0D, 0x12, 0x20)
C_ACCENT  = RGBColor(0xF0, 0xAB, 0x00)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT   = RGBColor(0xC8, 0xD6, 0xE5)
C_MUTED   = RGBColor(0x7A, 0x8B, 0xA0)
C_BOX     = RGBColor(0x1E, 0x30, 0x50)
C_BOX2    = RGBColor(0x24, 0x3B, 0x5E)
C_RED     = RGBColor(0xE0, 0x4B, 0x4B)
C_GREEN   = RGBColor(0x3E, 0xB4, 0x89)
C_BLUE    = RGBColor(0x4A, 0x9E, 0xE0)
C_PURPLE  = RGBColor(0x9B, 0x6D, 0xD0)
C_TEAL    = RGBColor(0x2E, 0xA8, 0x9D)
C_ORANGE  = RGBColor(0xE8, 0x85, 0x3D)

LEVEL_COLORS = {3: C_ACCENT, 2: C_BLUE, 1: C_TEAL, 0: C_PURPLE}


# ── Helpers ───────────────────────────────────────────────────────────────────

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def _bg(slide, color=None):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color or C_BG

def _rect(slide, left, top, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s

def _rrect(slide, left, top, w, h, fill_color, border_color=None):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill_color
    if border_color:
        s.line.color.rgb = border_color
        s.line.width = Pt(2)
    else:
        s.line.fill.background()
    return s

def _text(slide, text, left, top, w, h, size=20, bold=False, color=C_WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box

def _bullets(slide, items, left, top, w, h, size=22, color=C_LIGHT, bullet_color=C_ACCENT):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(10)
        p.alignment = PP_ALIGN.LEFT
        # Bullet marker
        r1 = p.add_run()
        r1.text = "  "
        r1.font.size = Pt(size)
        r1.font.color.rgb = bullet_color
        # Text
        r2 = p.add_run()
        r2.text = item
        r2.font.size = Pt(size)
        r2.font.color.rgb = color
    return box

def _notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text

def _img(slide, path, left, top, w, h):
    if path and Path(path).exists():
        slide.shapes.add_picture(str(path), left, top, w, h)

def _footer(slide, num):
    _rect(slide, Inches(0), Inches(7.25), SLIDE_W, Inches(0.04), C_ACCENT)
    _text(slide, f"{num} / {TOTAL}", Inches(12.0), Inches(7.28), Inches(1.2), Inches(0.22),
          size=9, color=C_MUTED, align=PP_ALIGN.RIGHT)
    _text(slide, "Baraghith (2025) \u00b7 From Prompts to Populations",
          Inches(0.3), Inches(7.28), Inches(5), Inches(0.22),
          size=9, color=C_MUTED)

def _title_bar(slide, title, subtitle=None):
    _text(slide, title, Inches(0.5), Inches(0.25), Inches(12), Inches(0.8),
          size=28, bold=True, color=C_WHITE)
    _rect(slide, Inches(0.5), Inches(1.05), Inches(6.0), Inches(0.04), C_ACCENT)
    if subtitle:
        _text(slide, subtitle, Inches(0.5), Inches(1.15), Inches(12), Inches(0.5),
              size=16, color=C_MUTED)

def _arrow_down(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()

def _arrow_up(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.UP_ARROW, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()

def _arrow_right(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()

# ── Image paths ───────────────────────────────────────────────────────────────
IMG_DIR = Path(__file__).parent.parent / "output" / "images"
SLUG = "from-prompts-to-populations"

def img(old_num, stype):
    return IMG_DIR / f"{SLUG}_s{old_num:02d}_{stype}.png"


# ==============================================================================
# SLIDE BUILDERS
# ==============================================================================

def slide_01_title(prs):
    """Title slide with atmospheric image."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _img(slide, img(1, "title"), Inches(7.4), Inches(0.3), Inches(5.6), Inches(6.9))
    _rect(slide, Inches(7.1), Inches(0.3), Inches(0.06), Inches(6.9), C_ACCENT)

    _text(slide, "From Prompts\nto Populations", Inches(0.5), Inches(1.2), Inches(6.4), Inches(2.5),
          size=40, bold=True)
    _rect(slide, Inches(0.5), Inches(3.8), Inches(4.0), Inches(0.06), C_ACCENT)
    _text(slide, "Cultural Evolution as a Meta-Frame\nfor AI Ontology and Individuation",
          Inches(0.5), Inches(4.1), Inches(6.4), Inches(1.5), size=20, color=C_LIGHT)
    _text(slide, "Baraghith (2025)", Inches(0.5), Inches(5.5), Inches(6), Inches(0.4),
          size=16, color=C_ACCENT)
    _notes(slide, "Based on Baraghith (2025). The paper asks: How many kinds of AI are out there? The answer depends on which criteria you use.")


def slide_02_practitioners(prs):
    """AI individuation depends on who you ask — 2x3 grid."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, "AI individuation depends on who you ask")

    practitioners = [
        ("Engineers", "architecture + weights", C_BLUE),
        ("RL Researchers", "behavioural policies", C_TEAL),
        ("Auditors", "instance traceability", C_GREEN),
        ("Roboticists", "embodiment + control", C_PURPLE),
        ("Governance", "purpose + context", C_ORANGE),
        ("Law", "accountable entity", C_RED),
    ]

    bw, bh = Inches(3.8), Inches(2.0)
    gap = Inches(0.3)
    start_x, start_y = Inches(0.5), Inches(1.7)

    for i, (who, what, color) in enumerate(practitioners):
        col, row = i % 3, i // 3
        x = start_x + col * (bw + gap)
        y = start_y + row * (bh + gap)
        _rrect(slide, x, y, bw, bh, C_BOX, color)
        _rect(slide, x + Inches(0.15), y + Inches(0.15), bw - Inches(0.3), Inches(0.05), color)
        _text(slide, who, x + Inches(0.2), y + Inches(0.3), bw - Inches(0.4), Inches(0.5),
              size=22, bold=True, color=color)
        _text(slide, what, x + Inches(0.2), y + Inches(0.9), bw - Inches(0.4), Inches(0.8),
              size=18, color=C_LIGHT)

    _text(slide, "Six criteria \u2014 six answers \u2014 frequent conflicts",
          Inches(0.5), Inches(6.3), Inches(12), Inches(0.5),
          size=18, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    _footer(slide, 2)
    _notes(slide, "Different communities already individuate AI differently. Engineers version by weights, RL researchers track policies, auditors need instance IDs, roboticists focus on embodiment, governance tracks purpose, law needs accountable entities.")


def slide_03_conflicts(prs):
    """These six criteria systematically conflict — conflict matrix."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, "These six criteria systematically conflict")

    conflicts = [
        ("Model", "Behaviour", "Same weights, different prompts\n\u2192 different behaviour", C_BLUE, C_TEAL),
        ("Model", "Instance ID", "One model \u2192 millions of sessions", C_BLUE, C_GREEN),
        ("Control-loop", "Model", "Same weights, different robot\n\u2192 different agent", C_PURPLE, C_BLUE),
        ("Instance ID", "Legal", "Millions of sessions \u2192 one\naccountable entity", C_GREEN, C_RED),
    ]

    bw = Inches(12.3)
    bh = Inches(1.0)
    start_x = Inches(0.5)
    start_y = Inches(1.7)

    for i, (left, right, example, lc, rc) in enumerate(conflicts):
        y = start_y + i * (bh + Inches(0.15))
        _rect(slide, start_x, y, bw, bh, C_BOX)
        _text(slide, left, start_x + Inches(0.2), y + Inches(0.1), Inches(2.0), Inches(0.4),
              size=20, bold=True, color=lc)
        _text(slide, "\u2260", start_x + Inches(2.2), y + Inches(0.1), Inches(0.5), Inches(0.4),
              size=24, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
        _text(slide, right, start_x + Inches(2.8), y + Inches(0.1), Inches(2.0), Inches(0.4),
              size=20, bold=True, color=rc)
        _text(slide, example, start_x + Inches(5.5), y + Inches(0.1), Inches(6.5), Inches(0.8),
              size=16, color=C_LIGHT)

    _text(slide, '"Each criterion targets a different level of organisation"',
          Inches(0.5), Inches(6.3), Inches(12), Inches(0.5),
          size=16, color=C_ACCENT, align=PP_ALIGN.CENTER)
    _footer(slide, 3)
    _notes(slide, "Same model + different prompts = same under model-instance but different under behavioural. Millions of sessions map to one legal entity. These conflicts are structural, not accidental.")


def slide_04_ontology_section(prs):
    """Section: The Ontology Problem."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _img(slide, img(3, "section"), Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    # Dark overlay at bottom
    s = _rect(slide, Inches(0), Inches(5.2), SLIDE_W, Inches(2.3), C_BG_DARK)
    s.fill.fore_color.rgb = C_BG_DARK
    _rect(slide, Inches(0), Inches(5.2), SLIDE_W, Inches(0.04), C_ACCENT)
    _text(slide, "The Ontology Problem", Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.0),
          size=34, bold=True)
    _text(slide, "Four philosophical positions on what AI \u201cis\u201d",
          Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.9), size=18, color=C_LIGHT)
    _footer(slide, 4)
    _notes(slide, "Before the evolutionary argument, the paper surveys four recent philosophical positions on AI ontology.")


def slide_05_four_views(prs):
    """Four positions on a 2D quadrant — individuation strength vs agency location."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, "Four positions span individuation and agency")

    # Axes
    cx, cy = Inches(6.5), Inches(4.2)
    # Horizontal axis
    _rect(slide, Inches(1.5), cy, Inches(10.5), Inches(0.04), C_MUTED)
    _text(slide, "Weak individuation", Inches(1.5), cy + Inches(0.15), Inches(3), Inches(0.3),
          size=14, color=C_MUTED)
    _text(slide, "Strong individuation", Inches(9.5), cy + Inches(0.15), Inches(2.5), Inches(0.3),
          size=14, color=C_MUTED, align=PP_ALIGN.RIGHT)
    # Vertical axis
    _rect(slide, cx, Inches(1.5), Inches(0.04), Inches(5.5), C_MUTED)
    _text(slide, "Agency in artifact", Inches(6.7), Inches(1.5), Inches(3), Inches(0.3),
          size=14, color=C_MUTED)
    _text(slide, "Agency in process", Inches(6.7), Inches(6.5), Inches(3), Inches(0.3),
          size=14, color=C_MUTED)

    # Quadrant positions (left=weak, right=strong; top=artifact, bottom=process)
    positions = [
        (Inches(8.5), Inches(2.0), "Ferrario (2025)", "Artifact Realism", C_BLUE),
        (Inches(2.0), Inches(2.0), "Cabitza et al.", "Cybork: ensembles", C_TEAL),
        (Inches(2.0), Inches(5.0), "Weinbaum & Veitas", "Process Ontology", C_PURPLE),
        (Inches(8.5), Inches(5.0), "Hawley (2019)", "Ontological Caution", C_ORANGE),
    ]
    for x, y, author, label, color in positions:
        _rrect(slide, x, y, Inches(3.2), Inches(1.2), C_BOX, color)
        _text(slide, author, x + Inches(0.15), y + Inches(0.1), Inches(2.9), Inches(0.4),
              size=14, color=C_MUTED)
        _text(slide, label, x + Inches(0.15), y + Inches(0.5), Inches(2.9), Inches(0.5),
              size=20, bold=True, color=color)

    _footer(slide, 5)
    _notes(slide, "These four positions differ on: how strongly they reify AI as individuals, and where they locate agency — in the artifact, in socio-technical practices, in open-ended processes, or in a moving conceptual frontier.")


def slide_06_pivot(prs):
    """No single criterion resolves AI identity — pivot to CET."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, "No single criterion resolves AI identity")

    # Quote block
    _rrect(slide, Inches(1.0), Inches(2.0), Inches(11.3), Inches(2.5), C_BOX, C_ACCENT)
    _text(slide, "\u201cThese conflicts cannot be resolved by choosing\none \u2018true\u2019 criterion \u2014 each targets a different\nlevel of organisation.\u201d",
          Inches(1.5), Inches(2.3), Inches(10.3), Inches(1.8),
          size=24, color=C_LIGHT, align=PP_ALIGN.CENTER)
    _text(slide, "\u2014 Baraghith (2025)", Inches(1.5), Inches(4.0), Inches(10.3), Inches(0.4),
          size=16, color=C_ACCENT, align=PP_ALIGN.RIGHT)

    # Arrow: Pluralism → Meta-frame needed
    _rrect(slide, Inches(2.5), Inches(5.3), Inches(3.0), Inches(0.8), C_BOX, C_MUTED)
    _text(slide, "Pluralism?", Inches(2.5), Inches(5.4), Inches(3.0), Inches(0.6),
          size=20, color=C_MUTED, align=PP_ALIGN.CENTER)
    _arrow_right(slide, Inches(5.7), Inches(5.45), Inches(1.5), Inches(0.5), C_ACCENT)
    _rrect(slide, Inches(7.5), Inches(5.3), Inches(3.5), Inches(0.8), C_BOX, C_ACCENT)
    _text(slide, "Meta-frame needed", Inches(7.5), Inches(5.4), Inches(3.5), Inches(0.6),
          size=20, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    _footer(slide, 6)
    _notes(slide, "Pluralism merely tolerates multiple criteria. But it cannot explain WHY these criteria exist, or WHY they conflict. A structural explanation is needed — this is where CET enters.")


def slide_07_cet_section(prs):
    """Section: Cultural Evolutionary Theory."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _img(slide, img(7, "section"), Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    s = _rect(slide, Inches(0), Inches(5.2), SLIDE_W, Inches(2.3), C_BG_DARK)
    _rect(slide, Inches(0), Inches(5.2), SLIDE_W, Inches(0.04), C_ACCENT)
    _text(slide, "Cultural Evolutionary Theory", Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.0),
          size=34, bold=True)
    _text(slide, "A unifying meta-frame, not just pluralism",
          Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.9), size=18, color=C_LIGHT)
    _footer(slide, 7)
    _notes(slide, "CET offers more than tolerance of multiple criteria. It provides a structural explanation for why they exist and how they relate.")


def slide_08_cet_vs_pluralism(prs):
    """CET explains what pluralism merely tolerates — 2-column comparison."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, "CET explains what pluralism merely tolerates")

    # Left column: Pluralism
    _rrect(slide, Inches(0.5), Inches(1.7), Inches(5.8), Inches(3.5), C_BOX, C_MUTED)
    _text(slide, "Methodological Pluralism", Inches(0.8), Inches(1.85), Inches(5.2), Inches(0.5),
          size=22, bold=True, color=C_MUTED)
    _rect(slide, Inches(0.8), Inches(2.35), Inches(2.5), Inches(0.04), C_MUTED)
    plur = ["\u2713  Tolerates multiple criteria", "\u2717  Cannot explain WHY they exist",
            "\u2717  No causal connections"]
    for i, t in enumerate(plur):
        c = C_GREEN if "\u2713" in t else C_RED
        _text(slide, t, Inches(0.8), Inches(2.8) + Inches(i * 0.7), Inches(5.2), Inches(0.5),
              size=20, color=c)

    # Right column: CET
    _rrect(slide, Inches(7.0), Inches(1.7), Inches(5.8), Inches(3.5), C_BOX, C_ACCENT)
    _text(slide, "Cultural Evolutionary Theory", Inches(7.3), Inches(1.85), Inches(5.2), Inches(0.5),
          size=22, bold=True, color=C_ACCENT)
    _rect(slide, Inches(7.3), Inches(2.35), Inches(3.2), Inches(0.04), C_ACCENT)
    cet = ["\u2713  Explains criteria via levels", "\u2713  Predicts which cluster",
           "\u2713  One coupled system"]
    for i, t in enumerate(cet):
        _text(slide, t, Inches(7.3), Inches(2.8) + Inches(i * 0.7), Inches(5.2), Inches(0.5),
              size=20, color=C_GREEN)

    # Arrow between columns
    _arrow_right(slide, Inches(6.35), Inches(3.8), Inches(0.5), Inches(0.4), C_ACCENT)

    _footer(slide, 8)
    _notes(slide, "Pluralism says 'use whichever criterion fits'. CET says 'each criterion targets a specific evolutionary level, and I can tell you which one and why.'")


def slide_09_vis(prs):
    """AI exhibits variation, inheritance, and selection — 3 columns."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, "AI development exhibits variation, inheritance, selection")

    pillars = [
        ("Variation", ["New architectures", "Fine-tuning"], C_BLUE),
        ("Inheritance", ["Model cloning", "API deployment"], C_TEAL),
        ("Selection", ["Benchmarks + RLHF", "Market adoption"], C_ACCENT),
    ]

    bw = Inches(3.6)
    gap = Inches(0.5)
    start_x = Inches(1.0)

    for i, (title, items, color) in enumerate(pillars):
        x = start_x + i * (bw + gap)
        _rrect(slide, x, Inches(1.7), bw, Inches(4.5), C_BOX, color)
        _rect(slide, x + Inches(0.15), Inches(1.85), bw - Inches(0.3), Inches(0.05), color)
        _text(slide, title, x + Inches(0.2), Inches(2.0), bw - Inches(0.4), Inches(0.5),
              size=24, bold=True, color=color, align=PP_ALIGN.CENTER)
        for j, item in enumerate(items):
            _text(slide, item, x + Inches(0.3), Inches(2.8) + Inches(j * 0.65), bw - Inches(0.6), Inches(0.5),
                  size=18, color=C_LIGHT, align=PP_ALIGN.CENTER)

    _text(slide, "The three conditions for cultural evolution are met",
          Inches(0.5), Inches(6.5), Inches(12), Inches(0.5),
          size=16, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    _footer(slide, 9)
    _notes(slide, "CET requires three conditions: variation (new architectures, fine-tunes), inheritance (cloning, forking, APIs), and selection (benchmarks, RLHF, adoption). AI development fulfils all three.")


def slide_10_hierarchy_section(prs):
    """Section: The Four-Level Hierarchy."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _img(slide, img(9, "section"), Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    s = _rect(slide, Inches(0), Inches(5.2), SLIDE_W, Inches(2.3), C_BG_DARK)
    _rect(slide, Inches(0), Inches(5.2), SLIDE_W, Inches(0.04), C_ACCENT)
    _text(slide, "The Four-Level Hierarchy", Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.0),
          size=34, bold=True)
    _text(slide, "Models \u2014 Systems \u2014 Episodes \u2014 Outputs",
          Inches(0.8), Inches(6.3), Inches(11.7), Inches(0.9), size=18, color=C_LIGHT)
    _footer(slide, 10)
    _notes(slide, "The core proposal: AI development operates across four nested levels, each with its own V/T/S dynamics.")


def slide_11_hierarchy(prs):
    """Four levels capture AI's evolutionary organisation — THE core diagram."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, "Four nested levels organise AI evolution")

    levels = [
        (3, "Model Lineages", "\u2248 Genetic lineages", "months\u2013years"),
        (2, "Configured Systems", "\u2248 Organisms", "days\u2013weeks"),
        (1, "Episodes & Instances", "\u2248 Interactions", "real-time"),
        (0, "Circulating Outputs", "\u2248 Cultural memes", "sec\u2013years"),
    ]

    bw = Inches(8.0)
    bh = Inches(1.05)
    gap = Inches(0.25)
    start_x = Inches(2.7)
    start_y = Inches(1.7)

    for i, (lvl, name, analogy, timescale) in enumerate(levels):
        y = start_y + i * (bh + gap)
        color = LEVEL_COLORS[lvl]
        _rrect(slide, start_x, y, bw, bh, C_BOX, color)
        _rect(slide, start_x, y, Inches(0.08), bh, color)
        _text(slide, f"Level {lvl}", start_x + Inches(0.2), y + Inches(0.1), Inches(1.2), Inches(0.4),
              size=20, bold=True, color=color)
        _text(slide, name, start_x + Inches(1.5), y + Inches(0.1), Inches(3.5), Inches(0.4),
              size=20, bold=True, color=C_WHITE)
        _text(slide, analogy, start_x + Inches(5.2), y + Inches(0.1), Inches(2.5), Inches(0.4),
              size=16, color=C_MUTED, align=PP_ALIGN.RIGHT)
        _text(slide, timescale, start_x + Inches(1.5), y + Inches(0.55), Inches(3.0), Inches(0.3),
              size=14, color=C_MUTED)

        if i < 3:
            _arrow_down(slide, start_x + Inches(3.8), y + bh + Inches(0.02), Inches(0.3), gap - Inches(0.04), LEVEL_COLORS[levels[i+1][0]])

    # Labels
    _text(slide, "HIGHER\nabstraction", Inches(0.5), Inches(2.0), Inches(2.0), Inches(0.8),
          size=14, color=C_MUTED, align=PP_ALIGN.CENTER)
    _text(slide, "LOWER\nabstraction", Inches(0.5), Inches(5.0), Inches(2.0), Inches(0.8),
          size=14, color=C_MUTED, align=PP_ALIGN.CENTER)
    _rect(slide, Inches(1.45), Inches(2.8), Inches(0.04), Inches(2.2), C_MUTED)
    _arrow_down(slide, Inches(1.32), Inches(4.8), Inches(0.3), Inches(0.3), C_MUTED)

    _footer(slide, 11)
    _notes(slide, "Level 3 = parameterised weight sets forming lineages. Level 2 = model + prompts + tools + guardrails. Level 1 = individual chat sessions or API calls. Level 0 = texts, images, code produced by AI.")


def slide_12_vts_table(prs):
    """Each level has its own V/T/S dynamics — compact table."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, "Each level has distinct evolutionary dynamics")

    # Headers
    headers = ["Level", "Variation", "Transmission", "Selection"]
    hw = [Inches(2.5), Inches(3.2), Inches(3.2), Inches(3.2)]
    hx = [Inches(0.5)]
    for i in range(1, 4):
        hx.append(hx[-1] + hw[i-1] + Inches(0.1))

    hy = Inches(1.7)
    rh = Inches(0.45)
    for i, (h, w, x) in enumerate(zip(headers, hw, hx)):
        _rect(slide, x, hy, w, rh, C_ACCENT)
        _text(slide, h, x + Inches(0.1), hy + Inches(0.05), w - Inches(0.2), rh,
              size=16, bold=True, color=C_BG_DARK, align=PP_ALIGN.CENTER)

    # Data
    rows = [
        (3, "L3: Models", "architectures", "cloning, distillation", "benchmarks"),
        (2, "L2: Systems", "tools, prompts", "templating, configs", "cost, regulation"),
        (1, "L1: Episodes", "context, seeds", "training data", "user feedback"),
        (0, "L0: Outputs", "sampling, edits", "sharing, scraping", "attention, ranking"),
    ]

    drh = Inches(0.7)
    gap = Inches(0.12)
    sy = hy + rh + Inches(0.1)

    for ri, (lvl, label, v, t, s) in enumerate(rows):
        y = sy + ri * (drh + gap)
        color = LEVEL_COLORS[lvl]
        for ci, (w, x) in enumerate(zip(hw, hx)):
            _rect(slide, x, y, w, drh, C_BOX)
        _rect(slide, hx[0], y, Inches(0.06), drh, color)
        _text(slide, label, hx[0] + Inches(0.15), y + Inches(0.15), hw[0] - Inches(0.3), drh - Inches(0.3),
              size=16, bold=True, color=color)
        _text(slide, v, hx[1] + Inches(0.15), y + Inches(0.15), hw[1] - Inches(0.3), drh - Inches(0.3),
              size=14, color=C_LIGHT)
        _text(slide, t, hx[2] + Inches(0.15), y + Inches(0.15), hw[2] - Inches(0.3), drh - Inches(0.3),
              size=14, color=C_LIGHT)
        _text(slide, s, hx[3] + Inches(0.15), y + Inches(0.15), hw[3] - Inches(0.3), drh - Inches(0.3),
              size=14, color=C_LIGHT)

    _footer(slide, 12)
    _notes(slide, "This table condenses the four level descriptions into one overview. Each level has its own variation, transmission, and selection mechanisms operating at different timescales.")


def slide_13_cross_level(prs):
    """Downward constrains, upward aggregates — cross-level causal pathways."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, "Downward causation constrains; upward aggregates")

    # 4 level boxes in center
    bw, bh = Inches(3.8), Inches(0.9)
    cx = Inches(4.8)
    start_y = Inches(1.8)
    gap = Inches(0.5)

    for i, (lvl, name) in enumerate([(3, "L3: Models"), (2, "L2: Systems"), (1, "L1: Episodes"), (0, "L0: Outputs")]):
        y = start_y + i * (bh + gap)
        color = LEVEL_COLORS[lvl]
        _rrect(slide, cx, y, bw, bh, C_BOX, color)
        _text(slide, name, cx + Inches(0.15), y + Inches(0.15), bw - Inches(0.3), Inches(0.5),
              size=20, bold=True, color=color, align=PP_ALIGN.CENTER)

    # Downward arrows (left)
    down_x = Inches(3.8)
    _rect(slide, down_x, start_y + bh, Inches(0.06), 3 * (bh + gap) - gap, C_RED)
    _text(slide, "CONSTRAINS \u2193", Inches(1.8), Inches(3.5), Inches(2.0), Inches(0.4),
          size=16, bold=True, color=C_RED, align=PP_ALIGN.CENTER)
    for i in range(3):
        ay = start_y + (i + 1) * (bh + gap) + Inches(0.2)
        _arrow_down(slide, down_x - Inches(0.1), ay, Inches(0.3), Inches(0.35), C_RED)

    # Upward arrows (right)
    up_x = Inches(9.3)
    _rect(slide, up_x, start_y + bh, Inches(0.06), 3 * (bh + gap) - gap, C_GREEN)
    _text(slide, "\u2191 AGGREGATES", Inches(9.6), Inches(3.5), Inches(2.0), Inches(0.4),
          size=16, bold=True, color=C_GREEN)
    for i in range(3):
        ay = start_y + i * (bh + gap) + bh + Inches(0.05)
        _arrow_up(slide, up_x - Inches(0.1), ay, Inches(0.3), Inches(0.35), C_GREEN)

    # Bypass L0→L3
    by_x = Inches(1.0)
    top_y = start_y + bh / 2
    bot_y = start_y + 3 * (bh + gap) + bh / 2
    _rect(slide, by_x, top_y, Inches(0.08), bot_y - top_y, C_ACCENT)
    _rect(slide, by_x, top_y, cx - by_x, Inches(0.04), C_ACCENT)
    _rect(slide, by_x, bot_y, cx - by_x, Inches(0.04), C_ACCENT)
    _arrow_up(slide, by_x - Inches(0.12), top_y - Inches(0.35), Inches(0.35), Inches(0.35), C_ACCENT)
    _text(slide, "L0 \u2192 L3\nDirect feedback", Inches(0.3), Inches(4.2), Inches(1.6), Inches(0.7),
          size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    _footer(slide, 13)
    _notes(slide, "Higher levels constrain lower-level variation. Lower levels aggregate into selection signals. The L0→L3 bypass (outputs re-entering training corpora) is a crucial direct feedback loop.")


def slide_14_mapping(prs):
    """The six criteria map onto different hierarchy levels."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, "Criteria map onto hierarchy levels, not against each other")

    # Level boxes on left, criteria on right
    levels_criteria = [
        (3, "Level 3: Models", ["Model-instance", "Behavioural/policy"], "Heritable information stored"),
        (2, "Level 2: Systems", ["Purpose", "Control-loop"], "Functional integration"),
        (1, "Level 1: Episodes", ["Instance ID"], "Selection events occur"),
        ("N", "Normative Overlay", ["Legal personhood"], "Applied to any level"),
    ]

    start_y = Inches(1.7)
    rh = Inches(1.1)
    gap = Inches(0.15)

    for i, (lvl, label, criteria, role) in enumerate(levels_criteria):
        y = start_y + i * (rh + gap)
        color = LEVEL_COLORS.get(lvl, C_MUTED)

        # Level box
        _rrect(slide, Inches(0.5), y, Inches(2.8), rh, C_BOX, color)
        _text(slide, label, Inches(0.7), y + Inches(0.15), Inches(2.4), Inches(0.7),
              size=18, bold=True, color=color)

        # Arrow
        _arrow_right(slide, Inches(3.4), y + Inches(0.3), Inches(0.4), Inches(0.4), color)

        # Criteria boxes
        cx = Inches(4.0)
        for j, crit in enumerate(criteria):
            _rrect(slide, cx + Inches(j * 3.0), y + Inches(0.1), Inches(2.8), rh - Inches(0.2), C_BOX2, color)
            _text(slide, crit, cx + Inches(j * 3.0 + 0.15), y + Inches(0.25), Inches(2.5), Inches(0.5),
                  size=16, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

        # Role description
        _text(slide, role, Inches(10.2), y + Inches(0.3), Inches(2.8), Inches(0.5),
              size=14, color=C_MUTED)

    _text(slide, "Conflicts resolved: criteria target different levels, not the same object",
          Inches(0.5), Inches(6.5), Inches(12), Inches(0.5),
          size=16, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    _footer(slide, 14)
    _notes(slide, "This is the resolution: the six criteria are not rivals. Model-instance and behavioural criteria track Level 3. Purpose and control-loop track Level 2. Instance ID tracks Level 1. Legal personhood is a normative overlay.")


def slide_15_rlhf(prs):
    """RLHF demonstrates multi-level dynamics — diamond cycle."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, "RLHF demonstrates multi-level evolutionary dynamics")

    # Diamond layout: L3 top, L2 right, L1 bottom, L0 left
    nw, nh = Inches(2.8), Inches(0.8)
    nodes = [
        (Inches(5.3), Inches(1.8), 3, "L3: Models"),
        (Inches(9.0), Inches(3.5), 2, "L2: Systems"),
        (Inches(5.3), Inches(5.2), 1, "L1: Episodes"),
        (Inches(1.5), Inches(3.5), 0, "L0: Outputs"),
    ]

    for x, y, lvl, label in nodes:
        color = LEVEL_COLORS[lvl]
        _rrect(slide, x, y, nw, nh, C_BOX, color)
        _text(slide, label, x + Inches(0.1), y + Inches(0.15), nw - Inches(0.2), Inches(0.5),
              size=20, bold=True, color=color, align=PP_ALIGN.CENTER)

    # Connecting lines (L-shaped, adjusted for 0.8" node height)
    # L3(right=8.1, mid_y=2.2) → L2(left=9.0, mid_y=3.9)
    _rect(slide, Inches(8.1), Inches(2.2), Inches(0.9), Inches(0.04), C_ACCENT)
    _rect(slide, Inches(8.96), Inches(2.2), Inches(0.04), Inches(1.7), C_ACCENT)
    # L2(bottom_mid_x=10.4, bottom=4.3) → L1(right=8.1, mid_y=5.6)
    _rect(slide, Inches(9.5), Inches(4.3), Inches(0.04), Inches(1.3), C_ACCENT)
    _rect(slide, Inches(8.1), Inches(5.56), Inches(1.44), Inches(0.04), C_ACCENT)
    # L1(left=5.3, mid_y=5.6) → L0(right=4.3, mid_y=3.9)
    _rect(slide, Inches(4.3), Inches(5.56), Inches(1.0), Inches(0.04), C_ACCENT)
    _rect(slide, Inches(4.3), Inches(3.9), Inches(0.04), Inches(1.7), C_ACCENT)
    # L0(top_mid_x=2.9, top=3.5) → L3(left=5.3, mid_y=2.2)
    _rect(slide, Inches(3.0), Inches(2.2), Inches(0.04), Inches(1.3), C_ACCENT)
    _rect(slide, Inches(3.0), Inches(2.2), Inches(2.3), Inches(0.04), C_ACCENT)

    _text(slide, "No single level owns the process",
          Inches(0.5), Inches(6.7), Inches(12), Inches(0.5),
          size=18, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    _footer(slide, 15)
    _notes(slide, "RLHF demonstrates multi-level dynamics: Level 1 episodes generate selection signals, Level 2 configuration filters which episodes occur, Level 3 weights get reshaped, Level 0 outputs circulate beyond the system.")


def slide_16_governance(prs):
    """Different governance interventions target different levels."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, "Governance interventions cascade through levels")

    interventions = [
        (3, "Training", "data + weights", "broadest", C_ACCENT),
        (2, "Systems", "oversight + audit", "targeted", C_BLUE),
        (1, "Episodes", "consent + feedback", "fine-grained", C_TEAL),
        (0, "Outputs", "watermarks + labels", "environmental", C_PURPLE),
    ]

    bw = Inches(2.6)
    bh = Inches(4.0)
    gap = Inches(0.5)
    start_x = Inches(0.5)
    start_y = Inches(1.7)

    for i, (lvl, title, what, scope, color) in enumerate(interventions):
        x = start_x + i * (bw + gap)
        _rrect(slide, x, start_y, bw, bh, C_BOX, color)
        _rect(slide, x + Inches(0.1), start_y + Inches(0.1), bw - Inches(0.2), Inches(0.05), color)
        _text(slide, f"Level {lvl}", x + Inches(0.15), start_y + Inches(0.2), bw - Inches(0.3), Inches(0.3),
              size=14, color=C_MUTED)
        _text(slide, title, x + Inches(0.15), start_y + Inches(0.5), bw - Inches(0.3), Inches(0.6),
              size=18, bold=True, color=color)
        _text(slide, what, x + Inches(0.15), start_y + Inches(1.3), bw - Inches(0.3), Inches(0.8),
              size=16, color=C_LIGHT)
        _text(slide, scope, x + Inches(0.15), start_y + bh - Inches(0.5), bw - Inches(0.3), Inches(0.4),
              size=16, bold=True, color=color, align=PP_ALIGN.CENTER)

        if i < 3:
            _arrow_right(slide, x + bw + Inches(0.08), start_y + bh / 2 - Inches(0.15), Inches(0.35), Inches(0.3), C_ACCENT)

    _text(slide, "Each intervention cascades differently through the coupled system",
          Inches(0.5), Inches(6.3), Inches(12), Inches(0.5),
          size=16, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    _footer(slide, 16)
    _notes(slide, "Mandatory watermarks (L0) reshape which system configs remain competitive. Human-in-the-loop (L2) constrains deployment designs. Restricting training data (L3) constrains the entire downstream tree.")


def slide_17_descriptive_normative(prs):
    """CET separates descriptive from normative questions — 2 columns."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, "CET separates descriptive from normative questions")

    # Left: Descriptive
    _rrect(slide, Inches(0.5), Inches(1.7), Inches(5.5), Inches(4.0), C_BOX, C_TEAL)
    _text(slide, "Descriptive (CET)", Inches(0.8), Inches(1.9), Inches(4.8), Inches(0.5),
          size=22, bold=True, color=C_TEAL)
    _rect(slide, Inches(0.8), Inches(2.4), Inches(2.5), Inches(0.04), C_TEAL)
    desc_items = ["How do AI systems vary?", "How are traits inherited?", "What gets selected, and why?"]
    for i, item in enumerate(desc_items):
        _text(slide, item, Inches(1.0), Inches(2.8) + Inches(i * 0.7), Inches(4.5), Inches(0.5),
              size=20, color=C_LIGHT)

    # Right: Normative
    _rrect(slide, Inches(7.3), Inches(1.7), Inches(5.5), Inches(4.0), C_BOX, C_PURPLE)
    _text(slide, "Normative (Ethics/Law)", Inches(7.6), Inches(1.9), Inches(4.8), Inches(0.5),
          size=22, bold=True, color=C_PURPLE)
    _rect(slide, Inches(7.6), Inches(2.4), Inches(3.0), Inches(0.04), C_PURPLE)
    norm_items = ["Which AI deserves moral status?", "Who bears responsibility?", "What should be regulated?"]
    for i, item in enumerate(norm_items):
        _text(slide, item, Inches(7.8), Inches(2.8) + Inches(i * 0.7), Inches(4.5), Inches(0.5),
              size=20, color=C_LIGHT)

    # Bridge
    _text(slide, "CET informs but does not determine",
          Inches(2.5), Inches(6.0), Inches(8), Inches(0.5),
          size=18, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
    _rect(slide, Inches(6.0), Inches(3.5), Inches(1.3), Inches(0.04), C_ACCENT)
    _text(slide, "informs \u2192", Inches(6.0), Inches(3.1), Inches(1.3), Inches(0.3),
          size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)

    _footer(slide, 17)
    _notes(slide, "Methodological strength: CET is independent of consciousness/sentience debates. It grounds individuation in operational mechanisms of copying, modification, and differential retention.")


def slide_18_future(prs):
    """Future work: fitness notions and reproduction predicates."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _title_bar(slide, "Three questions CET must still answer")

    items = [
        "Explicit reproduction predicates per level",
        "Level-specific fitness notions",
        "Pressure-testing on concrete governance cases",
    ]

    for i, item in enumerate(items):
        y = Inches(2.2) + Inches(i * 1.3)
        _rrect(slide, Inches(1.5), y, Inches(10.3), Inches(0.9), C_BOX, C_ACCENT)
        _text(slide, f"{i+1}", Inches(1.7), y + Inches(0.15), Inches(0.5), Inches(0.5),
              size=28, bold=True, color=C_ACCENT)
        _text(slide, item, Inches(2.5), y + Inches(0.2), Inches(9.0), Inches(0.5),
              size=22, color=C_LIGHT)

    _footer(slide, 18)
    _notes(slide, "The framework needs: explicit reproduction predicates at each level, level-specific fitness notions, and pressure-testing on concrete governance cases like the EU AI Act.")


def slide_19_closing(prs):
    """Closing slide with atmospheric image."""
    slide = _blank(prs)
    _bg(slide, C_BG)
    _img(slide, img(18, "closing"), Inches(0), Inches(0), Inches(5.8), SLIDE_H)
    _rect(slide, Inches(5.8), Inches(0.3), Inches(0.06), Inches(6.9), C_ACCENT)

    _text(slide, "One Process,\nMany Levels", Inches(6.3), Inches(1.5), Inches(6.5), Inches(2.5),
          size=36, bold=True)
    _rect(slide, Inches(6.3), Inches(3.8), Inches(3.0), Inches(0.06), C_ACCENT)
    _text(slide, "AI individuation is not a choice\nbetween criteria \u2014 it\u2019s recognising\nwhere each belongs in the\nevolutionary hierarchy.",
          Inches(6.3), Inches(4.1), Inches(6.5), Inches(2.0), size=20, color=C_LIGHT)

    _footer(slide, 19)
    _notes(slide, "Take-home: Rather than selecting one 'true' criterion for AI identity, CET provides a meta-frame where all six criteria are partial perspectives on a single multi-level evolutionary process.")


# ==============================================================================
# MAIN
# ==============================================================================

def main():
    project_dir = Path(__file__).parent.parent
    out_dir = project_dir / "output" / "presentations"
    out_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_practitioners(prs)
    slide_03_conflicts(prs)
    slide_04_ontology_section(prs)
    slide_05_four_views(prs)
    slide_06_pivot(prs)
    slide_07_cet_section(prs)
    slide_08_cet_vs_pluralism(prs)
    slide_09_vis(prs)
    slide_10_hierarchy_section(prs)
    slide_11_hierarchy(prs)
    slide_12_vts_table(prs)
    slide_13_cross_level(prs)
    slide_14_mapping(prs)
    slide_15_rlhf(prs)
    slide_16_governance(prs)
    slide_17_descriptive_normative(prs)
    slide_18_future(prs)
    slide_19_closing(prs)

    out_path = out_dir / "from-prompts-to-populations_v4.pptx"
    prs.save(str(out_path))

    print(json.dumps({
        "ok": True,
        "path": str(out_path),
        "slides": TOTAL,
        "version": 4,
    }, indent=2))


if __name__ == "__main__":
    main()
