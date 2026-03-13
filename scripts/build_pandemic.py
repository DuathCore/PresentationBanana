#!/usr/bin/env python3
"""
PresentationBanana — Pandemic/Infodemic Presentation Builder v1
===============================================================
Builds pandemic-infodemic_v1.pptx from generated diagram/chart/imagen PNGs.
Paper: Häusler & Baraghith (2023) — Biology & Philosophy 38:42
Style: dark-professional (Navy #121A2E + Gold #F0AB00)
"""
import json
import sys
from pathlib import Path

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print(json.dumps({"ok": False, "error": "pip install python-pptx"}))
    sys.exit(1)

# ── Dimensions ────────────────────────────────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL = 18

# ── Colors ────────────────────────────────────────────────────────────────────
C_BG      = RGBColor(0x12, 0x1A, 0x2E)
C_BG_DARK = RGBColor(0x0D, 0x12, 0x20)
C_ACCENT  = RGBColor(0xF0, 0xAB, 0x00)
C_WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT   = RGBColor(0xC8, 0xD6, 0xE5)
C_MUTED   = RGBColor(0x7A, 0x8B, 0xA0)
C_RED     = RGBColor(0xE0, 0x4B, 0x4B)

# ── Paths ─────────────────────────────────────────────────────────────────────
IMG_DIR = Path(__file__).parent.parent / "output" / "images"
SLUG = "pandemic-infodemic"


def img_diagram(n):
    return IMG_DIR / f"{SLUG}_s{n:02d}_diagram.png"


def img_chart(n):
    return IMG_DIR / f"{SLUG}_s{n:02d}_chart.png"


def img_imagen(n, stype):
    return IMG_DIR / f"{SLUG}_s{n:02d}_{stype}.png"


# ── Low-level helpers ─────────────────────────────────────────────────────────

def _blank(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide, color=None):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color or C_BG


def _rect(slide, left, top, w, h, color):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def _text(slide, text, left, top, w, h, size=20, bold=False,
          color=C_WHITE, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.color.rgb = color
    return box


def _notes(slide, text):
    if text:
        slide.notes_slide.notes_text_frame.text = text


def _img(slide, path, left, top, w, h):
    if path and Path(path).exists():
        slide.shapes.add_picture(str(path), left, top, w, h)


def _footer(slide, num):
    _rect(slide, Inches(0), Inches(7.25), SLIDE_W, Inches(0.04), C_ACCENT)
    _text(slide, f"{num} / {TOTAL}",
          Inches(12.0), Inches(7.28), Inches(1.2), Inches(0.22),
          size=9, color=C_MUTED, align=PP_ALIGN.RIGHT)
    _text(slide, "Häusler & Baraghith (2023) · Biology & Philosophy 38:42",
          Inches(0.3), Inches(7.28), Inches(6), Inches(0.22),
          size=9, color=C_MUTED)


def _title_bar(slide, title):
    _text(slide, title,
          Inches(0.5), Inches(0.18), Inches(12.3), Inches(0.9),
          size=26, bold=True, color=C_WHITE)
    _rect(slide, Inches(0.5), Inches(1.05), Inches(6.0), Inches(0.04), C_ACCENT)


def _diagram_slide(prs, num, title, diagram_path, notes_text):
    """Standard content slide: title bar + full-width diagram PNG."""
    slide = _blank(prs)
    _bg(slide)
    _title_bar(slide, title)
    _img(slide, diagram_path, Inches(0.3), Inches(1.2), Inches(12.7), Inches(5.9))
    _footer(slide, num)
    _notes(slide, notes_text)
    return slide


def _bullets_slide(prs, num, title, bullets, notes_text):
    """Text-only content slide with bullet points."""
    slide = _blank(prs)
    _bg(slide)
    _title_bar(slide, title)

    y = Inches(1.3)
    for bullet in bullets:
        box = slide.shapes.add_textbox(Inches(0.7), y, Inches(11.9), Inches(0.7))
        tf = box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = bullet
        r.font.size = Pt(20)
        r.font.color.rgb = C_LIGHT
        y += Inches(0.75)

    _footer(slide, num)
    _notes(slide, notes_text)
    return slide


# ══════════════════════════════════════════════════════════════════════════════
# SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def slide_01_title(prs):
    """Title slide: image right panel, text left."""
    slide = _blank(prs)
    _bg(slide)
    _img(slide, img_imagen(1, "title"),
         Inches(7.4), Inches(0.3), Inches(5.6), Inches(6.9))
    _rect(slide, Inches(7.1), Inches(0.3), Inches(0.06), Inches(6.9), C_ACCENT)

    _text(slide, "When Fake News Goes Viral",
          Inches(0.5), Inches(1.0), Inches(6.4), Inches(1.5),
          size=36, bold=True)
    _text(slide, "A Cultural Evolutionary Analysis",
          Inches(0.5), Inches(2.6), Inches(6.4), Inches(0.8),
          size=24, color=C_LIGHT)
    _rect(slide, Inches(0.5), Inches(3.55), Inches(4.5), Inches(0.06), C_ACCENT)
    _text(slide, "Häusler & Baraghith (2023)",
          Inches(0.5), Inches(3.75), Inches(6), Inches(0.5),
          size=18, color=C_ACCENT)
    _text(slide, "Biology & Philosophy 38:42",
          Inches(0.5), Inches(4.25), Inches(6), Inches(0.4),
          size=18, color=C_MUTED)

    _footer(slide, 1)
    _notes(slide, "This paper asks: why does misinformation spread like a biological virus? The authors use Cultural Evolutionary Theory — not just metaphor — to answer that question rigorously. Today we walk through their argument.")


def slide_02_roadmap(prs):
    _diagram_slide(prs, 2, "Four Steps: Problem → Theory → Mechanism → Solution",
                   img_diagram(2),
                   "The paper follows a tight four-part structure: first diagnosing the limits of the infodemic metaphor, then providing a theoretical fix via CET, unpacking the mechanisms, and finally proposing a solution. We will follow the same arc.")


def slide_03_infodemic(prs):
    """Text-only slide about WHO infodemic definition."""
    slide = _blank(prs)
    _bg(slide)
    _title_bar(slide, "WHO 2020: 'Infodemic' — But Is the Concept Grounded?")

    # Centered quote block
    _rect(slide, Inches(0.8), Inches(1.4), Inches(11.7), Inches(2.5), C_BG_DARK)
    _rect(slide, Inches(0.8), Inches(1.4), Inches(0.06), Inches(2.5), C_ACCENT)
    _text(slide, '"overabundance of information — some accurate, some not — that spreads during an epidemic"',
          Inches(1.2), Inches(1.6), Inches(11.0), Inches(1.2),
          size=22, color=C_WHITE)
    _text(slide, "— WHO (2020)",
          Inches(1.2), Inches(3.0), Inches(6), Inches(0.5),
          size=18, color=C_ACCENT)

    # Stats row — fixed: real numbers only, no definitions
    stats = [
        ("2020", "WHO coined\n'infodemic' term"),
        ("14,301", "Scientific papers\nused 'infodemic'"),
        ("COVID-19", "Primary driver\nof infodemic research"),
    ]
    for i, (val, lbl) in enumerate(stats):
        x = Inches(1.0 + i * 4.0)
        _rect(slide, x, Inches(4.2), Inches(3.5), Inches(1.8), C_BG_DARK)
        _rect(slide, x, Inches(4.2), Inches(3.5), Inches(0.05), C_ACCENT)
        _text(slide, val, x + Inches(0.2), Inches(4.4), Inches(3.1), Inches(0.7),
              size=28, bold=True, color=C_ACCENT)
        _text(slide, lbl, x + Inches(0.2), Inches(5.1), Inches(3.1), Inches(0.7),
              size=18, color=C_LIGHT)

    _text(slide, "Simon & Camargo (2021)",
          Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5),
          size=14, color=C_MUTED)

    _footer(slide, 3)
    _notes(slide, "The WHO coined 'infodemic' to describe the parallel spread of disinformation alongside COVID-19. The term exploded in scientific and media use — but was it ever properly grounded? That's the paper's opening problem.")


def slide_04_flaws(prs):
    _diagram_slide(prs, 4, "The Analogy Has Three Fatal Flaws",
                   img_diagram(4),
                   "Simon & Camargo (2021) identify three ways the virus-misinformation analogy breaks down. These aren't minor quibbles — they undermine the entire practical utility of calling something an infodemic. The authors take this critique seriously before trying to save the analogy.")


def slide_05_section(prs):
    """Section slide: full-bleed image with text overlay."""
    slide = _blank(prs)
    _bg(slide)
    _img(slide, img_imagen(5, "section"), Inches(0), Inches(0), SLIDE_W, SLIDE_H)
    _rect(slide, Inches(0), Inches(5.1), SLIDE_W, Inches(2.4), C_BG_DARK)
    _rect(slide, Inches(0), Inches(5.1), SLIDE_W, Inches(0.05), C_ACCENT)
    _text(slide, "CET Transforms Metaphor into Mechanism",
          Inches(0.8), Inches(5.25), Inches(11.7), Inches(1.1),
          size=34, bold=True)
    _text(slide, "Addressing all three shortcomings of the infodemic analogy",
          Inches(0.8), Inches(6.35), Inches(11.7), Inches(0.9),
          size=18, color=C_LIGHT)
    _footer(slide, 5)
    _notes(slide, "Rather than abandon the metaphor, the authors propose Cultural Evolutionary Theory as the theoretical scaffolding that gives it real explanatory power. This section introduces CET.")


def slide_06_cet(prs):
    _diagram_slide(prs, 6, "Culture Evolves Like Genes — CET in Brief",
                   img_diagram(6),
                   "CET maps biological evolutionary mechanisms onto cultural phenomena. Key insight for this paper: digital social media enables near-perfect fidelity copying — making large-scale diffusion of fake news observable and analyzable like genetic spread.")


def slide_07_biases(prs):
    _diagram_slide(prs, 7, "Cognitive Biases Drive Cultural Selection",
                   img_diagram(7),
                   "The heart of CET's explanatory contribution is the taxonomy of cognitive biases that make certain cultural information more likely to be copied. Fake news strategically exploits many of these simultaneously — which is why it spreads so effectively.")


def slide_08_chart(prs):
    _diagram_slide(prs, 8, "Negativity and Threat Bias — Fear Goes Viral",
                   img_chart(8),
                   "Two biases work together here: people preferentially believe and share negative information, and threat-related content is transmitted even more reliably than merely negative content. The pandemic's fear climate was perfect fuel for fake news virality.")


def slide_09_prestige(prs):
    _diagram_slide(prs, 9, "Prestige Bias — Politicians as Superspreaders",
                   img_diagram(9),
                   "Model-based prestige bias amplifies top-down misinformation disproportionately. A politician spreading fake news reaches vastly more people than regular users — and those interactions cascade. This is one of the most empirically documented mechanisms in the paper.")


def slide_10_mci_had(prs):
    """Text-only slide: MCI and HAD cognitive mechanisms."""
    slide = _blank(prs)
    _bg(slide)
    _title_bar(slide, "MCI and HAD — Conspiracy Theories Are Cognitively Sticky")

    # MCI box — drastically shortened to keywords only
    _rect(slide, Inches(0.5), Inches(1.3), Inches(5.9), Inches(4.5), C_BG_DARK)
    _rect(slide, Inches(0.5), Inches(1.3), Inches(5.9), Inches(0.06), C_ACCENT)
    _text(slide, "MCI — Minimally Counterintuitive",
          Inches(0.7), Inches(1.45), Inches(5.5), Inches(0.6),
          size=20, bold=True, color=C_ACCENT)
    _text(slide, "• Familiar base + 1 twist = sticky\n• Memorable + shareable by design",
          Inches(0.7), Inches(2.15), Inches(5.5), Inches(3.0),
          size=20, color=C_LIGHT)

    # HAD box — drastically shortened to keywords only
    _rect(slide, Inches(6.9), Inches(1.3), Inches(5.9), Inches(4.5), C_BG_DARK)
    _rect(slide, Inches(6.9), Inches(1.3), Inches(5.9), Inches(0.06), C_ACCENT)
    _text(slide, "HAD — Hidden Agency Detection",
          Inches(7.1), Inches(1.45), Inches(5.5), Inches(0.6),
          size=20, bold=True, color=C_ACCENT)
    _text(slide, "• Evolved bias: detect hidden agents\n• Random event → intentional actor",
          Inches(7.1), Inches(2.15), Inches(5.5), Inches(3.0),
          size=20, color=C_LIGHT)

    # Bottom insight — keyword form, no full sentence
    _rect(slide, Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.9), C_BG_DARK)
    _text(slide, "Conspiracy theories = cognitively optimized — exploit ancient algorithms",
          Inches(0.7), Inches(6.1), Inches(12.0), Inches(0.7),
          size=20, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)

    _footer(slide, 10)
    _notes(slide, "Conspiracy theories are not random — they are cognitively optimized. They use minimally counterintuitive structures (mostly plausible, one shocking twist) and exploit our evolved tendency to detect hidden agents. This is why the same structural templates recur across pandemics.")


def slide_11_trim(prs):
    _diagram_slide(prs, 11, "Filter Bubbles ARE TRIMs: Misinformation Gets Locked In",
                   img_diagram(11),
                   "Here CET makes its most original contribution. Filter bubbles don't just passively contain misinformation — they actively function as TRIMs: cultural barriers that prevent valid information from entering while maximizing the stickiness of false information within the bubble.")


def slide_12_quadrant(prs):
    _diagram_slide(prs, 12, "Fake News Evolves Toward Maximum Diffusion AND Retention",
                   img_diagram(12),
                   "This is the dynamic evolution of fake news within bubbles. Initially catchy variants win by diffusion potential. Over time, the bubble selects for variants with maximum retention potential — ones that are almost impossible to dislodge even with direct counter-evidence.")


def slide_13_maladaptive(prs):
    """Text-only slide: fake news as maladaptive trait."""
    slide = _blank(prs)
    _bg(slide)
    _title_bar(slide, "Fake News Harms Fitness — A Maladaptive Cultural Trait")

    # Central argument chain
    chain = [
        ("Fake news exposure", "C_BLUE"),
        ("Reduced vaccination intent", "C_ORANGE"),
        ("Health measure non-compliance", "C_ORANGE"),
        ("Increased infection risk", "C_RED"),
        ("Reduced biological fitness", "C_RED"),
    ]

    color_map = {
        "C_BLUE": RGBColor(0x4A, 0x9E, 0xE0),
        "C_ORANGE": RGBColor(0xE8, 0x85, 0x3D),
        "C_RED": RGBColor(0xE0, 0x4B, 0x4B),
    }

    for i, (label, c_key) in enumerate(chain):
        x = Inches(0.7 + i * 2.4)
        color = color_map[c_key]
        # Bug fix: only use add_shape (removed redundant _rect call)
        s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, Inches(2.0), Inches(2.1), Inches(1.5))
        s.fill.solid()
        s.fill.fore_color.rgb = RGBColor(0x1E, 0x30, 0x50)
        s.line.color.rgb = color
        s.line.width = Pt(1.5)
        _text(slide, label, x + Inches(0.1), Inches(2.3), Inches(1.9), Inches(0.9),
              size=18, color=C_WHITE, align=PP_ALIGN.CENTER)
        # Arrow (except last)
        if i < 4:
            _text(slide, "→", x + Inches(2.1), Inches(2.55), Inches(0.3), Inches(0.5),
                  size=22, color=C_ACCENT, align=PP_ALIGN.CENTER)

    # CET framing box — citation integrated
    _rect(slide, Inches(0.5), Inches(4.0), Inches(12.3), Inches(1.7), C_BG_DARK)
    _rect(slide, Inches(0.5), Inches(4.0), Inches(0.06), Inches(1.7), C_ACCENT)
    _text(slide, "CET framing: Fake news = maladaptive cultural variant\nAnalogy: mutation with fitness disadvantage → carrier harm (maladaptive)",
          Inches(0.8), Inches(4.15), Inches(11.8), Inches(1.0),
          size=20, color=C_LIGHT)
    _text(slide, "De Oliveira & Albuquerque (2021)",
          Inches(0.8), Inches(5.3), Inches(11.8), Inches(0.35),
          size=14, color=C_MUTED)

    _footer(slide, 13)
    _notes(slide, "CET frames fake news not just as false information but as a maladaptive cultural trait — one that, like a genetic mutation conferring disadvantage, reduces the biological fitness of its carriers. Anti-vax behavior is the clearest empirical case.")


def slide_14_pseudoscience(prs):
    _diagram_slide(prs, 14, "Pseudo-Science Immunizes Itself Against Refutation",
                   img_diagram(14),
                   "The most insidious mechanism: pseudo-scientific fake news has evolved epistemic defense mechanisms. Refutation attempts are reinterpreted as confirmation. This makes standard debunking — correcting false claims after the fact — structurally ineffective.")


def slide_15_debunking(prs):
    _diagram_slide(prs, 15, "Debunking Fails — The Case for Prebunking",
                   img_diagram(15),
                   "The empirical record on debunking is damning. Not only do warnings often fail to reach people — they can backfire, reinforcing conspiratorial distrust. Prebunking, by contrast, intervenes before beliefs calcify.")


def slide_16_inoculation(prs):
    _diagram_slide(prs, 16, "Psychological Inoculation — Mental Vaccination",
                   img_diagram(16),
                   "Psychological inoculation mirrors medical vaccination: expose individuals to a weakened version of misinformation tactics in a safe environment. The result is critical thinking skills that generalize — participants become better at detecting manipulation across contexts.")


def slide_17_unification(prs):
    _diagram_slide(prs, 17, "CET Makes the Metaphor Precise and Testable",
                   img_diagram(17),
                   "The paper's payoff: CET doesn't just support the metaphor — it transforms it into a testable theoretical framework. Both pandemics and infodemics can now be understood using the same evolutionary vocabulary, with precision about where the analogy holds and where it breaks down.")


def slide_18_closing(prs):
    """Closing slide: image left panel, text right."""
    slide = _blank(prs)
    _bg(slide)
    _img(slide, img_imagen(18, "closing"),
         Inches(0), Inches(0), Inches(5.8), SLIDE_H)
    _rect(slide, Inches(5.8), Inches(0.3), Inches(0.06), Inches(6.9), C_ACCENT)

    _text(slide, "Prebunking at Scale\n— The Path Forward",
          Inches(6.3), Inches(1.3), Inches(6.5), Inches(2.0),
          size=32, bold=True)
    _rect(slide, Inches(6.3), Inches(3.4), Inches(3.5), Inches(0.06), C_ACCENT)
    _text(slide, "CET tells us WHY fake news spreads,\nWHY it persists, and\nWHY debunking arrives too late.\n\nInoculation-based education is the\nevolutionarily informed strategy.",
          Inches(6.3), Inches(3.6), Inches(6.5), Inches(2.3),
          size=18, color=C_LIGHT)

    # Questions
    _text(slide, "Questions?",
          Inches(6.3), Inches(5.65), Inches(6.5), Inches(0.5),
          size=22, bold=True, color=C_ACCENT)

    # Citation strip
    _rect(slide, Inches(6.3), Inches(6.3), Inches(6.5), Inches(0.8), C_BG_DARK)
    _text(slide, "Häusler & Baraghith (2023) · Biology & Philosophy 38:42",
          Inches(6.5), Inches(6.42), Inches(6.0), Inches(0.4),
          size=14, color=C_MUTED)

    _footer(slide, 18)
    _notes(slide, "The takeaway: fighting the infodemic requires upstream intervention. CET tells us why fake news spreads, why it persists, and why debunking is too late. Inoculation-based education — especially before beliefs form — is the evolutionarily informed strategy.")


# ══════════════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════════════

def main():
    project_dir = Path(__file__).parent.parent
    out_dir = project_dir / "output" / "presentations"
    out_dir.mkdir(parents=True, exist_ok=True)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    slide_01_title(prs)
    slide_02_roadmap(prs)
    slide_03_infodemic(prs)
    slide_04_flaws(prs)
    slide_05_section(prs)
    slide_06_cet(prs)
    slide_07_biases(prs)
    slide_08_chart(prs)
    slide_09_prestige(prs)
    slide_10_mci_had(prs)
    slide_11_trim(prs)
    slide_12_quadrant(prs)
    slide_13_maladaptive(prs)
    slide_14_pseudoscience(prs)
    slide_15_debunking(prs)
    slide_16_inoculation(prs)
    slide_17_unification(prs)
    slide_18_closing(prs)

    out_path = out_dir / "pandemic-infodemic_v3.pptx"
    prs.save(str(out_path))

    print(json.dumps({
        "ok": True,
        "path": str(out_path),
        "slides": TOTAL,
        "version": 3,
        "slug": SLUG,
    }, indent=2))


if __name__ == "__main__":
    main()
